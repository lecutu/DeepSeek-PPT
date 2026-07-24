"""
python-pptx 桥接层：读取真实 PPTX → 引擎数据模型 → 写回 PPTX

职责：
  - 解析 pptx 的 shape 树 → SlideElement 列表
  - 推断元素角色（基于 shape 类型/位置/占位符）
  - 将修改后的元素坐标回写到 pptx
"""

from __future__ import annotations
from engine import (
    SlideElement, BBox, ContentRole, CollisionRole,
    EMU_PER_PT, DEFAULT_SLIDE_W, DEFAULT_SLIDE_H,
)
from typing import Optional
from pathlib import Path
import copy


def _emu_to_pt(emu: int) -> float:
    return emu / EMU_PER_PT


def _pt_to_emu(pt: float) -> int:
    return int(pt * EMU_PER_PT)


# ── PPTX → SlideElement ────────────────────────────────────
def parse_slide_to_elements(slide) -> list[SlideElement]:
    """
    Parse a python-pptx Slide object into Reflex Engine elements.
    Returns list of SlideElement with inferred roles.
    """
    elements = []
    for i, shape in enumerate(slide.shapes):
        # Skip non-visual or group shapes
        if shape.shape_type is None:
            continue

        elem_id = f"shape-{i:02d}"
        try:
            bbox = BBox(
                x=_emu_to_pt(shape.left),
                y=_emu_to_pt(shape.top),
                w=_emu_to_pt(shape.width),
                h=_emu_to_pt(shape.height),
            )
        except Exception:
            continue

        # ── Infer roles ────────────────────────────────────
        content_role = ContentRole.UNKNOWN
        collision_role = CollisionRole.FOREGROUND_CONTENT

        is_placeholder = hasattr(shape, 'is_placeholder') and shape.is_placeholder
        if is_placeholder:
            ph_type = shape.placeholder_format.type if hasattr(shape, 'placeholder_format') else None
            # python-pptx placeholder type constants: TITLE=1, BODY=2, SUBTITLE=3, etc.
            ph_type_val = int(ph_type) if ph_type is not None else 0

            if ph_type_val == 1:   # PP_PLACEHOLDER_TITLE → TITLE
                content_role = ContentRole.TITLE
            elif ph_type_val == 3:  # PP_PLACEHOLDER_SUBTITLE
                content_role = ContentRole.SUBTITLE
            elif ph_type_val == 2:  # PP_PLACEHOLDER_BODY
                content_role = ContentRole.BODY
            elif ph_type_val in (6, 7):  # PP_PLACEHOLDER_OBJECT, PP_PLACEHOLDER_PICTURE
                content_role = ContentRole.FIGURE

        # By shape type
        shape_type = str(shape.shape_type) if shape.shape_type else ""
        if "PICTURE" in shape_type:
            if content_role == ContentRole.UNKNOWN:
                content_role = ContentRole.FIGURE
        elif "TABLE" in shape_type:
            if content_role == ContentRole.UNKNOWN:
                content_role = ContentRole.BODY
        elif "GROUP" in shape_type:
            collision_role = CollisionRole.DECORATIVE

        # By position heuristics
        bbox_y_pct = bbox.y / DEFAULT_SLIDE_H
        bbox_area_pct = (bbox.w * bbox.h) / (DEFAULT_SLIDE_W * DEFAULT_SLIDE_H)

        # ── Extract text ───────────────────────────────────
        text = ""
        font_size_pt = 12.0
        font_explicit = False
        if shape.has_text_frame:
            text = shape.text_frame.text[:200]
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size_pt = run.font.size / 12700  # EMU → pt
                        font_explicit = True
                        break

        if content_role == ContentRole.UNKNOWN:
            if bbox_y_pct < 0.25 and bbox.w / DEFAULT_SLIDE_W >= 0.6:
                content_role = ContentRole.TITLE
            elif bbox_y_pct > 0.85:
                content_role = ContentRole.FOOTER
            elif bbox.w < 150 and bbox.h < 40:
                content_role = ContentRole.CAPTION
            elif bbox_area_pct > 0.85:
                content_role = ContentRole.BACKGROUND
                collision_role = CollisionRole.BACKGROUND_FILL
            elif shape.has_text_frame and text.strip():
                content_role = ContentRole.BODY

        # ── z-order ────────────────────────────────────────
        try:
            z_order = int(shape.xml.get("order", i))
        except Exception:
            z_order = i

        elem = SlideElement(
            id=elem_id,
            bbox=bbox,
            content_role=content_role,
            collision_role=collision_role,
            font_size_pt=font_size_pt,
            font_explicit=font_explicit,
            text=text,
            z_order=z_order,
        )
        elements.append(elem)

    return elements


# ── Write back ─────────────────────────────────────────────
def apply_element_positions(slide, elements: list[SlideElement]):
    """
    Write element bbox back to python-pptx shapes.
    Matches by index (shape-00 → slide.shapes[0]).
    """
    for elem in elements:
        idx = _extract_index(elem.id)
        if idx is None or idx >= len(slide.shapes):
            continue

        shape = slide.shapes[idx]
        try:
            shape.left = _pt_to_emu(elem.bbox.x)
            shape.top = _pt_to_emu(elem.bbox.y)
            shape.width = _pt_to_emu(elem.bbox.w)
            shape.height = _pt_to_emu(elem.bbox.h)
        except Exception:
            continue


def apply_single_element(slide, elem: SlideElement):
    """Write a single element's position back to the shape."""
    idx = _extract_index(elem.id)
    if idx is None or idx >= len(slide.shapes):
        return
    shape = slide.shapes[idx]
    try:
        shape.left = _pt_to_emu(elem.bbox.x)
        shape.top = _pt_to_emu(elem.bbox.y)
        shape.width = _pt_to_emu(elem.bbox.w)
        shape.height = _pt_to_emu(elem.bbox.h)
    except Exception:
        pass


def _extract_index(elem_id: str) -> Optional[int]:
    """Given 'shape-03', return 3."""
    try:
        return int(elem_id.split("-")[1])
    except (IndexError, ValueError):
        return None


# ── Presentation-level helpers ─────────────────────────────
def open_presentation(path: str):
    """Open pptx and return Presentation object."""
    from pptx import Presentation
    return Presentation(path)


def read_slide(prs, slide_idx: int) -> tuple:
    """Return (slide, elements) for given slide index."""
    slide = prs.slides[slide_idx]
    elements = parse_slide_to_elements(slide)
    return slide, elements


def save_presentation(prs, path: str):
    """Save presentation, creating backup if path exists."""
    import shutil
    target = Path(path)
    if target.exists():
        backup = target.with_suffix(".pptx.bak")
        shutil.copy2(target, backup)
    prs.save(str(target))
