"""
Host Adapter — Python 侧接口

抽象 COM/VSTO 通道，提供两个模式：
  - local  模式：python-pptx 文件级操作（Day 1 已验证，用于开发和测试）
  - com    模式：通过 Named Pipe 与 C# PPTReflexService 通信（Phase 2）

统一接口使得上层 reflex.py 和 mcp_server.py 不需要知道底层是哪种模式。
"""

from __future__ import annotations
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional, Protocol
import json
import struct
import time
import io
import os

# ═══════════════════════════════════════════════════════════
# DATA TYPES
# ═══════════════════════════════════════════════════════════

@dataclass
class ElementInfo:
    """Enriched element data — superset of what python-pptx can read."""
    id: str                    # "shape-00"
    left_pt: float
    top_pt: float
    width_pt: float
    height_pt: float
    role: str = "unknown"
    text: str = ""
    font_size_pt: float = 12.0
    font_explicit: bool = False
    is_placeholder: bool = False
    placeholder_type: int = 0     # 1=title, 2=body, 3=subtitle
    z_order: int = 0
    is_visible: bool = True
    is_grouped: bool = False

    # COM-only fields (None = unknown in local mode)
    text_overflow: Optional[bool] = None
    actual_font_size_pt: Optional[float] = None
    crop_left: Optional[float] = None   # 0.0–1.0 percentage
    crop_right: Optional[float] = None
    crop_top: Optional[float] = None
    crop_bottom: Optional[float] = None
    has_animation: Optional[bool] = None


@dataclass
class SlideInfo:
    index: int
    width_pt: float
    height_pt: float
    element_count: int
    layout_name: str = ""


@dataclass
class RenderedBounds:
    """Actual text rendering bounds — only available via COM."""
    text_width_pt: float
    text_height_pt: float
    box_width_pt: float
    box_height_pt: float
    overflow: bool
    overflow_lines: int = 0
    auto_shrunk_pt: Optional[float] = None  # font size after auto-fit


@dataclass
class StateChange:
    """Detected change in PowerPoint (human edit or external)."""
    elements_changed: list[str]
    slide_index: int
    timestamp: float


# ═══════════════════════════════════════════════════════════
# ABSTRACT INTERFACE (Protocol)
# ═══════════════════════════════════════════════════════════

class IHostAdapter(Protocol):
    """Interface that both LocalAdapter and ComAdapter implement."""

    # ── Connection ─────────────────────────────────────────
    def connect(self) -> bool: ...
    def disconnect(self): ...
    def ping(self) -> bool: ...

    # ── File ───────────────────────────────────────────────
    def open_presentation(self, path: str) -> dict: ...
    def save_presentation(self, path: str | None = None) -> dict: ...

    # ── Read ───────────────────────────────────────────────
    def read_slide(self, index: int) -> list[ElementInfo]: ...
    def read_rendered_bounds(self, elem_id: str) -> Optional[RenderedBounds]: ...

    # ── Write ──────────────────────────────────────────────
    def apply_positions(self, updates: list[dict]) -> dict: ...
    def apply_text(self, elem_id: str, text: str, font_pt: float | None = None) -> dict: ...
    def delete_element(self, elem_id: str) -> dict: ...

    # ── Render ─────────────────────────────────────────────
    def render_slide_png(self, slide_index: int, dpi: int = 200) -> bytes | None: ...

    # ── Events ─────────────────────────────────────────────
    def poll_state_change(self, slide_index: int) -> Optional[StateChange]: ...


# ═══════════════════════════════════════════════════════════
# LOCAL ADAPTER (python-pptx)
# ═══════════════════════════════════════════════════════════

class LocalAdapter:
    """
    Uses python-pptx for file-level operations.
    No real-time PowerPoint interaction — reads/writes files.
    COM-only fields are always None.
    """

    def __init__(self):
        self._prs = None
        self._path: str = ""
        self._last_elements: dict[int, dict[str, tuple[float, float, float, float]]] = {}

    def connect(self) -> bool:
        """Local mode always succeeds — no external process needed."""
        return True

    def disconnect(self):
        if self._prs:
            self._prs = None

    def ping(self) -> bool:
        return self._prs is not None

    # ── File ───────────────────────────────────────────────

    def open_presentation(self, path: str) -> dict:
        from pptx import Presentation
        if not Path(path).exists():
            raise FileNotFoundError(f"File not found: {path}")
        self._prs = Presentation(path)
        self._path = path
        w_emu = self._prs.slide_width or 12192000  # default 16:9
        h_emu = self._prs.slide_height or 6858000
        return {
            "slides": len(self._prs.slides),
            "slide_width_pt": w_emu / 12700,
            "slide_height_pt": h_emu / 12700,
            "mode": "local",
        }

    def save_presentation(self, path: str | None = None) -> dict:
        target = path or self._path
        self._prs.save(target)
        return {"saved_to": target, "mode": "local"}

    # ── Read ───────────────────────────────────────────────

    def read_slide(self, index: int) -> list[ElementInfo]:
        from engine import EMU_PER_PT
        slide = self._prs.slides[index]
        elements = []

        for i, shape in enumerate(slide.shapes):
            if shape.shape_type is None:
                continue
            try:
                left = shape.left / EMU_PER_PT
                top = shape.top / EMU_PER_PT
                w = shape.width / EMU_PER_PT
                h = shape.height / EMU_PER_PT
            except Exception:
                continue

            text = ""
            font_pt = 12.0
            font_explicit = False
            if shape.has_text_frame:
                text = shape.text_frame.text[:200]
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_pt = run.font.size / 12700
                            font_explicit = True
                            break

            is_ph = hasattr(shape, 'is_placeholder') and shape.is_placeholder
            ph_type = 0
            if is_ph and hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                ph_type = int(shape.placeholder_format.type) if shape.placeholder_format.type else 0

            ei = ElementInfo(
                id=f"shape-{i:02d}",
                left_pt=left, top_pt=top,
                width_pt=w, height_pt=h,
                text=text,
                font_size_pt=font_pt,
                font_explicit=font_explicit,
                is_placeholder=is_ph,
                placeholder_type=ph_type,
                z_order=i,
            )
            elements.append(ei)

        # Snapshot for change detection
        self._last_elements[index] = {
            e.id: (e.left_pt, e.top_pt, e.width_pt, e.height_pt)
            for e in elements
        }
        return elements

    def read_rendered_bounds(self, elem_id: str) -> Optional[RenderedBounds]:
        """Not available in local mode — requires COM."""
        return None

    # ── Write ──────────────────────────────────────────────

    def apply_positions(self, updates: list[dict]) -> dict:
        """updates: [{"id": "shape-05", "left_pt": 36, "top_pt": 140, ...}, ...]"""
        from engine import EMU_PER_PT
        applied = 0
        for update in updates:
            eid = update["id"]
            idx = _shape_index(eid)
            if idx is None or idx >= len(self._prs.slides[0].shapes):
                continue
            shape = self._prs.slides[0].shapes[idx]
            try:
                shape.left = int(update["left_pt"] * EMU_PER_PT)
                shape.top = int(update["top_pt"] * EMU_PER_PT)
                shape.width = int(update["width_pt"] * EMU_PER_PT)
                shape.height = int(update["height_pt"] * EMU_PER_PT)
                applied += 1
            except Exception:
                continue
        return {"applied": applied, "total": len(updates), "mode": "local"}

    def apply_text(self, elem_id: str, text: str, font_pt: float | None = None) -> dict:
        idx = _shape_index(elem_id)
        if idx is None or idx >= len(self._prs.slides[0].shapes):
            return {"error": f"Element {elem_id} not found"}
        shape = self._prs.slides[0].shapes[idx]
        if shape.has_text_frame:
            shape.text_frame.text = text
            if font_pt and shape.text_frame.paragraphs:
                from pptx.util import Pt
                shape.text_frame.paragraphs[0].runs[0].font.size = Pt(font_pt)
        return {"ok": True}

    def delete_element(self, elem_id: str) -> dict:
        idx = _shape_index(elem_id)
        if idx is None or idx >= len(self._prs.slides[0].shapes):
            return {"error": f"Element {elem_id} not found"}
        # python-pptx doesn't easily delete shapes; mark by setting size to 0
        shape = self._prs.slides[0].shapes[idx]
        shape.width = 0
        shape.height = 0
        return {"deleted": elem_id, "mode": "local"}

    # ── Render ─────────────────────────────────────────────

    def render_slide_png(self, slide_index: int, dpi: int = 200) -> bytes | None:
        """Not available in local mode without LibreOffice."""
        return None

    # ── Events ─────────────────────────────────────────────

    def poll_state_change(self, slide_index: int) -> Optional[StateChange]:
        """Compare current slide against last snapshot."""
        elements = self.read_slide(slide_index)
        changed = []
        old_snapshot = self._last_elements.get(slide_index, {})
        for e in elements:
            if e.id in old_snapshot:
                old = old_snapshot[e.id]
                if (abs(e.left_pt - old[0]) > 0.5 or
                    abs(e.top_pt - old[1]) > 0.5 or
                    abs(e.width_pt - old[2]) > 0.5 or
                    abs(e.height_pt - old[3]) > 0.5):
                    changed.append(e.id)
        if changed:
            return StateChange(
                elements_changed=changed,
                slide_index=slide_index,
                timestamp=time.time(),
            )
        return None


# ═══════════════════════════════════════════════════════════
# COM ADAPTER (Named Pipe → C# PPTReflexService)
# ═══════════════════════════════════════════════════════════

class ComAdapter:
    """
    Communicates with a local C# service via Windows Named Pipe.
    The C# service hosts the PowerPoint COM automation.
    """

    DEFAULT_PIPE = r"\\.\pipe\ppt_reflex_pipe"
    RECONNECT_MAX = 3
    TIMEOUT_MS = 5000
    RENDER_TIMEOUT_MS = 30000

    def __init__(self, pipe_name: str = DEFAULT_PIPE):
        self._pipe_name = pipe_name
        self._pipe = None
        self._req_id = 0
        self._connected = False

    def connect(self) -> bool:
        attempts = 0
        while attempts < self.RECONNECT_MAX:
            try:
                # Named pipe client (Windows-only)
                self._pipe = open(self._pipe_name, 'r+b', buffering=0)
                self._connected = True
                return True
            except (FileNotFoundError, OSError):
                attempts += 1
                if attempts < self.RECONNECT_MAX:
                    time.sleep(1.0)
        return False

    def disconnect(self):
        if self._pipe:
            try:
                self._pipe.close()
            except Exception:
                pass
            self._pipe = None
            self._connected = False

    def ping(self) -> bool:
        if not self._connected or not self._pipe:
            return False
        try:
            resp = self._request("ping", {})
            return resp.get("status") == "pong"
        except Exception:
            return False

    # ── Internal ───────────────────────────────────────────

    def _request(self, method: str, params: dict, binary_data: bytes | None = None,
                 timeout_ms: int | None = None) -> dict:
        """Send a JSON-RPC-style request over the pipe and read response."""
        if not self._pipe:
            raise ConnectionError("Not connected to PowerPoint service")

        self._req_id += 1
        msg = json.dumps({
            "id": f"req-{self._req_id:04d}",
            "method": method,
            "params": params,
        }, ensure_ascii=False).encode('utf-8')

        # Write: 4-byte length prefix + JSON payload [+ binary payload]
        header = struct.pack('<I', len(msg))
        self._pipe.write(header + msg)
        if binary_data:
            self._pipe.write(binary_data)
        self._pipe.flush()

        # Read response: 4-byte length prefix + JSON
        timeout = timeout_ms or self.TIMEOUT_MS
        start = time.time()
        while True:
            if self._pipe.readable():
                break
            if time.time() - start > timeout / 1000:
                raise TimeoutError(f"No response from service within {timeout}ms")
            time.sleep(0.01)

        length_bytes = self._pipe.read(4)
        if len(length_bytes) < 4:
            raise ConnectionError("Pipe closed while reading response")
        msg_len = struct.unpack('<I', length_bytes)[0]
        response_data = self._pipe.read(msg_len)
        return json.loads(response_data.decode('utf-8'))

    # ── File ───────────────────────────────────────────────

    def open_presentation(self, path: str) -> dict:
        resp = self._request("open_presentation", {"path": path})
        if resp.get("status") != "ok":
            raise RuntimeError(f"Failed to open: {resp.get('message')}")
        return resp["result"]

    def save_presentation(self, path: str | None = None) -> dict:
        return self._request("save_presentation", {"path": path}).get("result", {})

    # ── Read ───────────────────────────────────────────────

    def read_slide(self, index: int) -> list[ElementInfo]:
        resp = self._request("read_elements", {"slide_idx": index})
        return [
            ElementInfo(
                id=e["id"],
                left_pt=e["left_pt"], top_pt=e["top_pt"],
                width_pt=e["width_pt"], height_pt=e["height_pt"],
                text=e.get("text", ""),
                font_size_pt=e.get("font_size_pt", 12.0),
                font_explicit=e.get("font_explicit", False),
                is_placeholder=e.get("is_placeholder", False),
                placeholder_type=e.get("placeholder_type", 0),
                z_order=e.get("z_order", 0),
                is_visible=e.get("is_visible", True),
                text_overflow=e.get("text_overflow"),
                actual_font_size_pt=e.get("actual_font_size_pt"),
                crop_left=e.get("crop_left"),
                crop_right=e.get("crop_right"),
                crop_top=e.get("crop_top"),
                crop_bottom=e.get("crop_bottom"),
            )
            for e in resp["elements"]
        ]

    def read_rendered_bounds(self, elem_id: str) -> Optional[RenderedBounds]:
        resp = self._request("rendered_bounds", {"element_id": elem_id})
        if resp.get("status") != "ok":
            return None
        rb = resp["result"]
        return RenderedBounds(**rb)

    # ── Write ──────────────────────────────────────────────

    def apply_positions(self, updates: list[dict]) -> dict:
        resp = self._request("apply_positions", {"updates": updates})
        return resp.get("result", {})

    def apply_text(self, elem_id: str, text: str, font_pt: float | None = None) -> dict:
        params = {"element_id": elem_id, "text": text}
        if font_pt is not None:
            params["font_pt"] = font_pt
        return self._request("apply_text", params).get("result", {})

    def delete_element(self, elem_id: str) -> dict:
        return self._request("delete_element", {"element_id": elem_id}).get("result", {})

    # ── Render ─────────────────────────────────────────────

    def render_slide_png(self, slide_index: int, dpi: int = 200) -> bytes | None:
        resp = self._request("render_slide_png",
                             {"slide_index": slide_index, "dpi": dpi},
                             timeout_ms=self.RENDER_TIMEOUT_MS)
        if resp.get("status") != "ok":
            return None
        # PNG data is sent as binary after the JSON response
        remaining = self._pipe.read()
        return remaining if remaining else None

    # ── Events ─────────────────────────────────────────────

    def poll_state_change(self, slide_index: int) -> Optional[StateChange]:
        resp = self._request("poll_state_change", {"slide_index": slide_index})
        if resp.get("status") != "ok" or not resp.get("result", {}).get("changed"):
            return None
        result = resp["result"]
        return StateChange(
            elements_changed=result["elements_changed"],
            slide_index=slide_index,
            timestamp=time.time(),
        )


# ═══════════════════════════════════════════════════════════
# ADAPTER FACTORY
# ═══════════════════════════════════════════════════════════

def create_adapter(mode: str = "local", pipe_name: str = "") -> IHostAdapter:
    """
    Factory: returns the correct adapter based on mode.

    mode:
      "local"  → python-pptx (always available, file-level)
      "com"    → Named Pipe → C# PPTReflexService (Windows + PowerPoint required)
    """
    if mode == "com":
        adapter = ComAdapter(pipe_name or ComAdapter.DEFAULT_PIPE)
        if not adapter.connect():
            raise ConnectionError(
                "Cannot connect to PowerPoint service. "
                "Ensure PPTReflexService.exe is running.\n"
                "Falling back to local mode."
            )
        return adapter
    return LocalAdapter()


# ── helper ─────────────────────────────────────────────────

def _shape_index(elem_id: str) -> int | None:
    """Parse 'shape-NN' → int index."""
    try:
        return int(elem_id.split("-")[1])
    except (IndexError, ValueError):
        return None
