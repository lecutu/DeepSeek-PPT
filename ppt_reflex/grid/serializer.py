"""
grid/serializer.py — Grid 状态 ↔ PPT 文件（唯一碰 python-pptx 的地方）

读:  PPT 文件 → InformationGrid（重建信息层）
写:  InformationGrid → PPT 文件（原子写入）
分类: python-pptx shape → ContentType
"""

from __future__ import annotations
import copy

from .types import GridConfig, ContentType, InfoCell, ElementPayload
from .info_grid import InformationGrid
from .positioning import bbox_to_fine_cells


# ═══════════════════════════════════════════════════════════
# READ: PPT → InformationGrid
# ═══════════════════════════════════════════════════════════

def ppt_to_grid(ppt_path: str, slide_index: int,
                config: GridConfig | None = None) -> InformationGrid:
    """从 PPT 文件读取指定幻灯片，重建信息层。

    Args:
        ppt_path: .pptx 文件路径
        slide_index: 0-indexed 幻灯片序号
        config: 网格配置

    Returns:
        填充好的 InformationGrid
    """
    from pptx import Presentation

    cfg = config or GridConfig()
    grid = InformationGrid(cfg)
    prs = Presentation(ppt_path)
    slides = list(prs.slides)

    if slide_index >= len(slides):
        raise IndexError(f"Slide {slide_index} not found (total: {len(slides)})")

    slide = slides[slide_index]
    slide_w_emu = prs.slide_width    # EMU
    slide_h_emu = prs.slide_height

    for shape in slide.shapes:
        try:
            left_pt = shape.left / 12700   # EMU → pt (python-pptx)
            top_pt = shape.top / 12700
            width_pt = shape.width / 12700
            height_pt = shape.height / 12700
        except Exception:
            continue

        content_type = classify_shape(shape)
        shape_id = _shape_id(shape)
        fine_cells = bbox_to_fine_cells(left_pt, top_pt, width_pt, height_pt, cfg)
        if not fine_cells:
            continue

        is_template = _is_template_element(shape)

        grid.occupy(
            fine_cells,
            owner_id=shape_id,
            content_type=content_type,
            z_order=_z_order(shape),
            locked=is_template,
            source="template" if is_template else "human",
        )

    return grid


# ═══════════════════════════════════════════════════════════
# WRITE: InformationGrid → PPT
# ═══════════════════════════════════════════════════════════

def grid_to_ppt(grid: InformationGrid, config: GridConfig, ppt_path: str,
               phase1_rects: dict | None = None,
               phase1_payloads: dict | None = None,
               decorations: list[dict] | None = None) -> None:
    """渲染管线：Phase 1 pt 直传（首选）→ info_grid 回退（旧架构）→ 装饰层（后渲染）。

    phase1_rects: {elem_id: (x, y, w, h)} — Phase 1 精确 pt，渲染首选
    phase1_payloads: {elem_id: (ContentType, ElementPayload)} — 渲染内容直传
    decorations:  Phase 2 装饰列表
    """
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    ALIGN_MAP = {
        "LEFT": PP_ALIGN.LEFT,
        "CENTER": PP_ALIGN.CENTER,
        "RIGHT": PP_ALIGN.RIGHT,
    }

    occ = grid.all_occupied()
    rects = phase1_rects or {}
    payloads = phase1_payloads or {}
    decos = decorations or []

    if not occ and not rects and not decos:
        return

    prs = Presentation()
    # 强制 16:9 宽屏，与 GridConfig canvas_w_pt=960/canvas_h_pt=540 对齐
    from pptx.util import Pt as _Pt
    prs.slide_width = _Pt(config.canvas_w_pt)
    prs.slide_height = _Pt(config.canvas_h_pt)
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # ── 信息层（先渲染 = 底层）──
    rendered_ids: set[str] = set()

    # 优先：Phase 1 精确 pt rects + payloads
    for eid, (x, y, w, h) in rects.items():
        rendered_ids.add(eid)
        x, y, w, h = _clamp_bbox(x, y, w, h, config)
        ct, payload = payloads.get(eid, (None, None))
        if ct is None:
            ct = _get_type_for_id(eid, occ.get(eid, set()), grid)
        if payload is None and eid in occ:
            cell_sample = grid.get_cell(next(iter(occ.get(eid, set())), ""))
            if cell_sample:
                payload = cell_sample.payload
        lock_flag = False
        if eid in occ:
            cs = grid.get_cell(next(iter(occ.get(eid, set())), ""))
            if cs:
                lock_flag = cs.locked
        if payload is not None and not lock_flag:
            if ct == ContentType.IMAGE and payload.image_path:
                _render_image(slide, x, y, w, h, payload)
                continue
            if ct == ContentType.CONNECTOR:
                _render_connector(slide, payload, grid.config)
                continue
            _render_payload(slide, x, y, w, h, ct, payload, ALIGN_MAP)
            continue

        _render_fallback(slide, x, y, w, h, ct)

    # 回退：info_grid 中未被 Phase 1 rects 覆盖的元素（旧架构兼容）
    for owner_id, fine_cells in occ.items():
        if owner_id in rendered_ids:
            continue
        bbox = _cells_union(fine_cells, grid.config)
        if bbox is None:
            continue
        x, y, w, h = bbox
        x, y, w, h = _clamp_bbox(x, y, w, h, config)
        content_type = _get_type_for_id(owner_id, fine_cells, grid)
        cell_sample = grid.get_cell(next(iter(fine_cells)))
        locked = cell_sample.locked if cell_sample else False
        payload = cell_sample.payload if cell_sample else None

        if payload is not None and not locked:
            if content_type == ContentType.IMAGE and payload.image_path:
                _render_image(slide, x, y, w, h, payload)
                continue
            if content_type == ContentType.CONNECTOR:
                _render_connector(slide, payload, grid.config)
                continue
            _render_payload(slide, x, y, w, h, content_type, payload, ALIGN_MAP)
            continue

        _render_fallback(slide, x, y, w, h, content_type)

    # ── 装饰层（后渲染 = 顶层，不被信息层遮挡）──
    for d in decos:
        if d.get("deco_type") == "arrow" and d.get("x2"):
            _render_connector_abs(slide, d, config, ALIGN_MAP)

    prs.save(ppt_path)


def _render_fallback(slide, x, y, w, h, content_type):
    """空占位形状（无 payload 时）。"""
    from pptx.util import Pt
    if content_type == ContentType.TEXT:
        box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        box.text_frame.text = ""
    elif content_type == ContentType.TEXTBOX:
        box = slide.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(h))
        box.text_frame.text = ""
    elif content_type == ContentType.IMAGE:
        slide.shapes.add_picture(_placeholder_png(), Pt(x), Pt(y), Pt(w), Pt(h))
    elif content_type == ContentType.SHAPE:
        slide.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(h))
    else:
        box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        box.text_frame.text = ""


# ═══════════════════════════════════════════════════════════
# BBOX CLAMP — 硬裁剪，防止任何算错的形状出页
# ═══════════════════════════════════════════════════════════

def _clamp_bbox(x: float, y: float, w: float, h: float,
                config: GridConfig) -> tuple[float, float, float, float]:
    """Clamp bbox to slide bounds. Guarantees no shape exceeds canvas."""
    SW = config.canvas_w_pt
    SH = config.canvas_h_pt
    M = 0.0  # safe margin already applied upstream; hard-clamp here
    cx = max(M, min(x, SW - 1))
    cy = max(M, min(y, SH - 1))
    cw = max(1, min(w, SW - cx))
    ch = max(1, min(h, SH - cy))
    return cx, cy, cw, ch


# ═══════════════════════════════════════════════════════════
# SHAPE ID MAP — ElementPayload.shape_id → python-pptx MSO_SHAPE
# ═══════════════════════════════════════════════════════════

_SHAPE_MAP = {
    "rectangle":            1,
    "rounded_rectangle":    5,
    "diamond":              4,
    "ellipse":              9,
    "chevron":             55,
    "pentagon":            56,
    "hexagon":              9,
    "star5":               92,
    "star8":               93,
    "triangle":             7,
    "right_triangle":       8,
    "right_arrow":         33,
    "left_arrow":          34,
    "up_arrow":            35,
    "down_arrow":          36,
    "striped_right_arrow": 93,
    "arc":                 19,
    "moon":                20,
    "parallelogram":        7,
    "trapezoid":            8,
    "plus":                11,
    "bevel":               15,
    "can":                 22,
    "cube":                16,
    "donut":               23,
    "lightning_bolt":      22,
    "heart":               21,
    "cloud":               179,
    "banner":              68,
    "seal5":               110,
    "seal8":               111,
    "flowchart_process":    1,    # rectangle
    "flowchart_decision":   4,    # diamond
    "flowchart_data":       6,    # parallelogram
    "flowchart_predefined": 1,    # rectangle
    "flowchart_internal_storage": 1,
    "flowchart_document":   1,
    "flowchart_multidocument": 1,
    "flowchart_terminator": 5,    # rounded rectangle
    "flowchart_preparation": 6,   # hexagon-like
    "flowchart_manual_input": 6,
    "flowchart_manual_operation": 6,
    "flowchart_connector":   9,   # circle
    "flowchart_offpage_connector": 6,
    "flowchart_card":        1,
    "flowchart_punched_tape": 1,
    "flowchart_summing_junction": 9,
    "flowchart_or":           9,
    "flowchart_collate":      9,
    "flowchart_sort":         9,
    "flowchart_extract":      9,
    "flowchart_merge":        7,
    "flowchart_offline_storage": 1,
    "flowchart_online_storage": 22,
    "flowchart_magnetic_tape": 1,
    "flowchart_magnetic_disk": 22,
    "flowchart_magnetic_drum": 22,
    "flowchart_display":      1,
    "flowchart_delay":        1,
}


def _lookup_shape(shape_id: str) -> int:
    """Map human-readable shape_id to python-pptx MSO_SHAPE integer.

    Returns 1 (rectangle) for unknown keys.
    """
    return _SHAPE_MAP.get(shape_id.lower(), 1)


# ═══════════════════════════════════════════════════════════
# RENDER: CONNECTOR — with anchor support
# ═══════════════════════════════════════════════════════════

_ANCHOR_FACTORS = {
    "top":      (0.5, 0.0),
    "bottom":   (0.5, 1.0),
    "left":     (0.0, 0.5),
    "right":    (1.0, 0.5),
    "center":   (0.5, 0.5),
}


def _anchor_point(col: int, row: int, anchor: str, cw: float) -> tuple[float, float]:
    """Compute anchor coordinate in pt from fine-cell address + anchor keyword."""
    fx, fy = _ANCHOR_FACTORS.get(anchor, (0.5, 0.5))
    return (col + fx) * cw, (row + fy) * cw


def _render_connector(slide, payload: "ElementPayload", config: "GridConfig") -> None:
    """Render an arrow connector between two grid cells.

    Uses connector_from / connector_to cell addresses → anchor point coordinates.
    Draws a straight line with arrowhead at target end.

    Anchor keywords: "top"|"bottom"|"left"|"right"|"center"

    Priority: _abs_x1/_abs_y1 pt coords (Phase 2 direct) > cell-address derivation.
    """
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    from .positioning import parse_cell

    # Phase 2 direct pt coords — bypass cell discretization
    if payload._abs_x2 or payload._abs_y2:
        sx, sy = payload._abs_x1, payload._abs_y1
        ex, ey = payload._abs_x2, payload._abs_y2
        sx = max(0, min(sx, config.canvas_w_pt - 1))
        sy = max(0, min(sy, config.canvas_h_pt - 1))
        ex = max(0, min(ex, config.canvas_w_pt - 1))
        ey = max(0, min(ey, config.canvas_h_pt - 1))
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, Pt(sx), Pt(sy), Pt(ex), Pt(ey))
        connector.line.color.rgb = RGBColor(*payload.line_color)
        connector.line.width = Pt(payload.line_width_pt)
        _set_arrowhead(connector)
        return

    fc = payload.connector_from.strip()
    tc = payload.connector_to.strip()
    if not fc or not tc:
        return

    try:
        src_col, src_row = parse_cell(fc)
        dst_col, dst_row = parse_cell(tc)
    except ValueError:
        return

    cw = config.fine_cell_pt

    sx, sy = _anchor_point(src_col, src_row, payload.connector_anchor_from, cw)
    ex, ey = _anchor_point(dst_col, dst_row, payload.connector_anchor_to, cw)

    # Clamp endpoints to slide bounds
    sx = max(0, min(sx, config.canvas_w_pt - 1))
    sy = max(0, min(sy, config.canvas_h_pt - 1))
    ex = max(0, min(ex, config.canvas_w_pt - 1))
    ey = max(0, min(ey, config.canvas_h_pt - 1))

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, Pt(sx), Pt(sy), Pt(ex), Pt(ey))
    connector.line.color.rgb = RGBColor(*payload.line_color)
    connector.line.width = Pt(payload.line_width_pt)

    _set_arrowhead(connector)


def _set_arrowhead(connector) -> None:
    """Set triangle arrowhead at target end."""
    from lxml import etree
    from pptx.oxml.ns import qn
    spPr = connector._element.find(qn('p:spPr'))
    if spPr is not None:
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            tail = ln.find(qn('a:tailEnd'))
            if tail is None:
                tail = etree.SubElement(ln, qn('a:tailEnd'))
            tail.set('type', 'triangle')
            tail.set('w', 'med')
            tail.set('len', 'med')


def _render_connector_abs(slide, d: dict, config: "GridConfig", align_map: dict) -> None:
    """Render Phase 2 decoration arrow with absolute pt coords + optional text label."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE

    x1 = max(0, min(d["x1"], config.canvas_w_pt - 1))
    y1 = max(0, min(d["y1"], config.canvas_h_pt - 1))
    x2 = max(0, min(d["x2"], config.canvas_w_pt - 1))
    y2 = max(0, min(d["y2"], config.canvas_h_pt - 1))

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, Pt(x1), Pt(y1), Pt(x2), Pt(y2))
    connector.line.color.rgb = RGBColor(*d["line_color"])
    connector.line.width = Pt(d["line_width_pt"])
    _set_arrowhead(connector)

    # ── Text label at midpoint ──
    if d.get("text", "").strip():
        mx = (x1 + x2) / 2
        my = (y1 + y2) / 2
        lw = max(36, d.get("font_size", 9) * len(d["text"].replace("\n", "")) * 0.65)
        lh = max(18, d.get("font_size", 9) * d["text"].count("\n") * 1.5 + d.get("font_size", 9) * 2)
        lx = max(0, min(mx - lw / 2, config.canvas_w_pt - lw))
        ly = max(0, min(my - lh / 2, config.canvas_h_pt - lh))
        box = slide.shapes.add_textbox(Pt(lx), Pt(ly), Pt(lw), Pt(lh))
        box.text_frame.word_wrap = True
        tf = box.text_frame
        for i, line in enumerate(d["text"].split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.font.name = d.get("font_name", "Calibri")
            para.font.size = Pt(d.get("font_size", 9))
            para.font.color.rgb = RGBColor(*d["font_color"])
            para.alignment = align_map.get("CENTER", 1)
# ═══════════════════════════════════════════════════════════
# RENDER: TEXT / TEXTBOX (with shape_id)
# ═══════════════════════════════════════════════════════════

def _render_payload(slide, x: float, y: float, w: float, h: float,
                    content_type: ContentType, p: ElementPayload,
                    align_map: dict) -> None:
    """Render a single element with full payload content."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    fc = RGBColor(*p.font_color) if p.font_color else RGBColor(0x22, 0x22, 0x44)
    has_fill = p.fill_color is not None
    is_code = (p.font_name in ("Consolas", "Courier New", "Source Code Pro", "Roboto Mono")
               and has_fill)

    # ── Base shape ──
    shape_type = _lookup_shape(p.shape_id) if p.shape_id else 1  # default = rectangle
    if is_code or has_fill or p.shape_id:
        shape = slide.shapes.add_shape(shape_type, Pt(x), Pt(y), Pt(w), Pt(h))
        if has_fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*p.fill_color)
        shape.line.fill.background()
        tf = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        tf = shape.text_frame

    tf.word_wrap = True
    tf.auto_size = None  # no auto-shrink; overflow caught by engine pre-check

    # Margins
    pad = Pt(12) if is_code else Pt(6)
    tf.margin_left = pad
    tf.margin_right = pad
    tf.margin_top = pad
    tf.margin_bottom = pad

    # ── Text content ──
    text = p.text.strip()
    if not text:
        # 有框无文 → 占位标记，防静默白框
        text = "(no text)"
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(max(8, p.font_size * 0.8))
        run.font.color.rgb = RGBColor(0xCC, 0x33, 0x33)  # 红色警告
        run.font.italic = True
        return

    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.name = p.font_name
        para.font.size = Pt(p.font_size)
        para.font.color.rgb = fc
        para.font.bold = p.font_bold
        para.alignment = align_map.get(p.alignment, align_map["LEFT"])
        para.line_spacing = p.line_spacing
        para.space_after = Pt(0)
        para.space_before = Pt(0)


# ═══════════════════════════════════════════════════════════
# RENDER: IMAGE
# ═══════════════════════════════════════════════════════════

def _render_image(slide, x: float, y: float, w: float, h: float,
                  p: ElementPayload) -> None:
    """Render an image element with fit/fill/crop_center mode."""
    import os
    from pptx.util import Pt, Emu
    from PIL import Image

    if not p.image_path or not os.path.isfile(p.image_path):
        return

    img = Image.open(p.image_path)
    img_w_px, img_h_px = img.size
    if img_w_px == 0 or img_h_px == 0:
        return

    target_w_pt = w
    target_h_pt = h
    img_aspect = img_w_px / img_h_px
    box_aspect = target_w_pt / target_h_pt

    if p.fit_mode == "fill":
        # Stretch to fill box exactly (distorts if aspect mismatch)
        pic = slide.shapes.add_picture(p.image_path, Pt(x), Pt(y), Pt(target_w_pt), Pt(target_h_pt))

    elif p.fit_mode == "crop_center":
        # Scale to cover box, then crop overflow (centered)
        if img_aspect > box_aspect:
            # Image wider → scale to match height, crop left/right
            scale = target_h_pt / img_h_px
            scaled_w = img_w_px * scale
            crop_x = x - (scaled_w - target_w_pt) / 2
            pic = slide.shapes.add_picture(p.image_path, Pt(crop_x), Pt(y), Pt(scaled_w), Pt(target_h_pt))
        else:
            # Image taller → scale to match width, crop top/bottom
            scale = target_w_pt / img_w_px
            scaled_h = img_h_px * scale
            crop_y = y - (scaled_h - target_h_pt) / 2
            pic = slide.shapes.add_picture(p.image_path, Pt(x), Pt(crop_y), Pt(target_w_pt), Pt(scaled_h))

    else:  # "fit" (default)
        # Scale to fit inside box, keep aspect ratio, centered
        if img_aspect > box_aspect:
            # Image wider → constrained by width
            fit_w = target_w_pt
            fit_h = target_w_pt / img_aspect
        else:
            # Image taller → constrained by height
            fit_h = target_h_pt
            fit_w = target_h_pt * img_aspect
        fit_x = x + (target_w_pt - fit_w) / 2
        fit_y = y + (target_h_pt - fit_h) / 2
        pic = slide.shapes.add_picture(p.image_path, Pt(fit_x), Pt(fit_y), Pt(fit_w), Pt(fit_h))


# ═══════════════════════════════════════════════════════════
# CLASSIFY: shape → ContentType
# ═══════════════════════════════════════════════════════════

def classify_shape(shape) -> ContentType:
    """python-pptx shape → ContentType 分类。

    逻辑:
      PICTURE       → IMAGE
      TABLE         → TABLE
      CHART         → CHART
      LINE/CONNECTOR → CONNECTOR
      TEXT_BOX      → 有 fill 且无文字 → SHAPE
                     → 有 fill 且有文字 → TEXTBOX
                     → 无 fill 有文字 → TEXT
      其他          → UNKNOWN
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        shape_type = shape.shape_type
    except Exception:
        return ContentType.UNKNOWN

    # Picture
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return ContentType.IMAGE

    # Table
    if shape_type == MSO_SHAPE_TYPE.TABLE:
        return ContentType.TABLE

    # Chart
    if shape_type == MSO_SHAPE_TYPE.CHART:
        return ContentType.CHART

    # Line / Connector
    if shape_type in (MSO_SHAPE_TYPE.LINE,):
        return ContentType.CONNECTOR

    # Placeholder / Text Box / Auto Shape
    if shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER,
                      MSO_SHAPE_TYPE.AUTO_SHAPE):
        from pptx.enum.dml import MSO_FILL_TYPE
        has_text = False
        has_fill = False

        try:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    has_text = True
        except Exception:
            pass

        try:
            fill = shape.fill
            ft = fill.type
            # BACKGROUND/NONE → no explicit fill. SOLID/PATTERN/GRADIENT/etc → has fill.
            if ft is not None and ft not in (MSO_FILL_TYPE.BACKGROUND,):
                has_fill = True
        except Exception:
            pass

        if has_text and has_fill:
            return ContentType.TEXTBOX
        if has_text:
            return ContentType.TEXT
        if has_fill:
            return ContentType.SHAPE
        return ContentType.UNKNOWN

    if shape_type == MSO_SHAPE_TYPE.GROUP:
        # 不深入组内，整个组标记为 UNKNOWN
        return ContentType.UNKNOWN

    return ContentType.UNKNOWN


# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════

def _shape_id(shape) -> str:
    """提取 shape 的唯一稳定 ID。"""
    try:
        return f"shape-{shape.shape_id}"
    except Exception:
        try:
            return f"shape-{hash(shape.name)}"
        except Exception:
            return f"shape-{id(shape)}"


def _z_order(shape) -> int:
    """shape 的 z-order。"""
    try:
        return shape.z_order
    except Exception:
        return 0


def _is_template_element(shape) -> bool:
    """判断 shape 是否来自母版/版式（非用户添加）。"""
    try:
        if getattr(shape, 'is_placeholder', False):
            return True
    except Exception:
        pass
    try:
        # python-pptx: 检查是否在 slide layout 上
        if hasattr(shape, 'part') and 'slideLayout' in str(type(shape.part)):
            return True
    except Exception:
        pass
    return False


def _cells_union(fine_cells: set[str], config: GridConfig) -> tuple | None:
    """一组信息格地址 → 最小包围矩形 (x, y, w, h) pt。"""
    from .positioning import parse_cell, cell_name

    if not fine_cells:
        return None
    parsed = [parse_cell(c) for c in fine_cells]
    min_col = min(p[0] for p in parsed)
    max_col = max(p[0] for p in parsed)
    min_row = min(p[1] for p in parsed)
    max_row = max(p[1] for p in parsed)
    x = min_col * config.fine_cell_pt
    y = min_row * config.fine_cell_pt
    w = (max_col - min_col + 1) * config.fine_cell_pt
    h = (max_row - min_row + 1) * config.fine_cell_pt
    return (x, y, w, h)


def _get_type_for_id(owner_id: str, cells: set[str], grid: InformationGrid) -> ContentType:
    """从信息格中取该 owner 的类型（取第一个非空 cell 的 content_type）。"""
    for addr in sorted(cells):
        cell = grid.get_cell(addr)
        if cell and cell.content_type:
            return cell.content_type
    return ContentType.UNKNOWN


def _placeholder_png() -> str:
    """返回一个 1×1 透明的 PNG 作为占位图。"""
    import tempfile, base64, os
    # 最小的 1x1 透明 PNG
    png_b64 = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk"
        "+P+/HgAFhAJ/qlOqBAAAAABJRU5ErkJggg=="
    )
    data = base64.b64decode(png_b64)
    tmp = os.path.join(tempfile.gettempdir(), "_ppt_reflex_placeholder.png")
    if not os.path.exists(tmp):
        with open(tmp, 'wb') as f:
            f.write(data)
    return tmp
