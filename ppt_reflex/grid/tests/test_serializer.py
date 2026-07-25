"""grid/tests/test_serializer.py — Grid ↔ PPT 往返一致性"""
import sys, os
from ppt_reflex.grid.types import GridConfig, ContentType
from ppt_reflex.grid.serializer import classify_shape, ppt_to_grid, grid_to_ppt


def test_classify_shape_types():
    import pptx
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Text box with text
    tb = slide.shapes.add_textbox(Pt(36), Pt(36), Pt(200), Pt(50))
    tb.text_frame.text = "Hello"
    assert classify_shape(tb) == ContentType.TEXT
    print("[PASS] classify: TEXT")

    # Shape with fill (rectangle)
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(36), Pt(100), Pt(200), Pt(50))
    shape.fill.solid()
    shape.fill.fore_color.rgb = pptx.dml.color.RGBColor(0xEE, 0xEE, 0xEE)
    # RECTANGLE is AUTO_SHAPE, has fill, no text → SHAPE
    assert classify_shape(shape) == ContentType.SHAPE
    print("[PASS] classify: SHAPE (fill, no text)")

    # Shape with fill + text
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(36), Pt(160), Pt(200), Pt(50))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = pptx.dml.color.RGBColor(0xCC, 0xCC, 0xCC)
    shape2.text_frame.text = "Colored box"
    assert classify_shape(shape2) == ContentType.TEXTBOX
    print("[PASS] classify: TEXTBOX (fill + text)")

    os.remove("_temp_test.pptx") if os.path.exists("_temp_test.pptx") else None
    prs.save("_temp_test.pptx")

    # Read back
    grid = ppt_to_grid("_temp_test.pptx", 0)
    occ = grid.all_occupied()
    assert len(occ) == 3, f"Expected 3 elements, got {len(occ)}"
    print(f"[PASS] ppt_to_grid: {len(occ)} elements read back")

    os.remove("_temp_test.pptx")


def test_round_trip():
    """写入→读取→一致性检查"""
    import tempfile
    from ppt_reflex.grid.info_grid import InformationGrid

    config = GridConfig()
    grid_in = InformationGrid(config)
    grid_in.occupy_bbox(36, 36, 200, 50, "s01", ContentType.TEXT, source="agent")
    grid_in.occupy_bbox(36, 100, 200, 50, "s02", ContentType.IMAGE, source="agent")
    grid_in.occupy_bbox(300, 36, 200, 200, "s03", ContentType.TEXTBOX, source="human")

    tmp = os.path.join(tempfile.gettempdir(), "_ppt_reflex_roundtrip.pptx")
    grid_to_ppt(grid_in, config, tmp)
    assert os.path.exists(tmp), "PPT file not created"
    print(f"[PASS] grid_to_ppt: {tmp} ({os.path.getsize(tmp)} bytes)")

    grid_out = ppt_to_grid(tmp, 0, config)
    occ_out = grid_out.all_occupied()
    assert len(occ_out) == 3, f"Round-trip lost elements: {len(occ_out)} vs 3"
    print(f"[PASS] ppt_to_grid: {len(occ_out)} elements recovered")

    # density check
    d_in = grid_in.density()
    d_out = grid_out.density()
    print(f"  density: {d_in:.3f} → {d_out:.3f}")

    os.remove(tmp)


if __name__ == "__main__":
    test_classify_shape_types()
    test_round_trip()
    print("\n✓ All serializer tests PASSED")
