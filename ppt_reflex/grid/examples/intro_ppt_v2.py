"""
intro_ppt_v2.py — 好看一点的 PPT Reflex Grid 介绍
"""
import sys, os, tempfile
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

OUT = os.path.join(tempfile.gettempdir(), "PPT_Reflex_Grid_Intro_v2.pptx")
prs = Presentation()
prs.slide_width  = Emu(960 * 12700)
prs.slide_height = Emu(540 * 12700)
blank = prs.slide_layouts[6]

# ── colors ────────────────────────────────────────────
BG   = RGBColor(0xFF, 0xFF, 0xFF)
BGL  = RGBColor(0xF5, 0xF5, 0xFA)
AC1  = RGBColor(0x2D, 0x5B, 0xD7)  # primary blue
AC2  = RGBColor(0x00, 0xA8, 0x8F)  # teal accent
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LGRAY = RGBColor(0xE0, 0xE4, 0xEB)
RED  = RGBColor(0xE0, 0x4F, 0x5F)

# ── helpers ───────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(blank)

def bg(slide, c=BG):
    r = slide.shapes.add_shape(1, 0, 0, Pt(960), Pt(540))
    r.fill.solid(); r.fill.fore_color.rgb = c
    r.line.fill.background()

def rect(slide, x, y, w, h, fill=None, radius=None, line_color=None):
    s = slide.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*fill) if isinstance(fill, tuple) else fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = RGBColor(*line_color)
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if radius:
        s.shadow.inherit = False
    return s

def text_box(slide, x, y, w, h, txt, size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = align; p.font.name = font_name
    return tb

def multi_text(slide, x, y, w, lines, leading=1.3):
    """lines = [(text, size, color, bold, align), ...]"""
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(400))
    tf = tb.text_frame; tf.word_wrap = True
    for i, parts in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        txt, size, color, bold, align = parts[0], parts[1] if len(parts)>1 else 14, \
            RGBColor(*parts[2]) if len(parts)>2 and isinstance(parts[2],tuple) else (parts[2] if len(parts)>2 else DARK), \
            parts[3] if len(parts)>3 else False, parts[4] if len(parts)>4 else PP_ALIGN.LEFT
        p.text = txt; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.bold = bold; p.alignment = align; p.font.name = "Microsoft YaHei"
        p.space_after = Pt(2)
    return tb

def card(slide, x, y, w, h, title, body, accent=AC1):
    rect(slide, x, y, w, h, fill=BGL)
    rect(slide, x, y, w, 4, fill=accent)
    text_box(slide, x+20, y+16, w-40, 30, title, size=15, bold=True)
    text_box(slide, x+20, y+48, w-40, h-60, body, size=12, color=GRAY)

def hline(slide, x, y, w):
    rect(slide, x, y, w, 1, fill=LGRAY)


# ═══════════════════════════════════════════════════════════
# COVER
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, 0, 0, 960, 540, fill=DARK)
rect(s, 0, 400, 960, 6, fill=AC1)

multi_text(s, 80, 140, 800, [
    ("PPT Reflex Engine", 42, (255,255,255), True, PP_ALIGN.LEFT),
    ("grid/ — 感知型网格画布", 24, (0xA0,0xC8,0xFF), False, PP_ALIGN.LEFT),
])
text_box(s, 80, 440, 800, 40, "两层定位 · 事前拦截 · 交互矩阵 · 分层供给 · 零破坏", size=14, color=RGBColor(0x88,0x90,0xA0))

# ═══════════════════════════════════════════════════════════
# PROBLEM — dark section
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK)
rect(s, 0, 0, 960, 540, fill=DARK)

text_box(s, 80, 60, 800, 40, "旧架构：5 个痛点", size=28, color=RGBColor(0xFF,0xFF,0xFF), bold=True)
hline(s, 80, 110, 200)

problems = [
    ("Agent 盲操作", "算坐标→写PPT→audit→发现问题→改→再审计\nPPT 被污染好几轮才勉强对"),
    ("事后检测", "audit 是[体检],不是[门禁]\nAgent 不知道放这里会不会撞, 放完才知道"),
    ("持久化bug", "跨 slide 切换时修复只留在内存\n保存后发现旧问题又回来了"),
    ("字体重叠死角", "字体/间距/密度能检测但无法修复\nwarn_only 策略 - 看到问题说[不管]"),
    ("token 浪费", "每 slide 425 tokens 全量审计\nAgent 理解的是坐标数字，不是空间概念"),
]
for i, (t, d) in enumerate(problems):
    x = 60 + (i % 3) * 290
    y = 160 + (i // 3) * 170
    card(s, x, y, 260, 150, t, d, accent=RED if i < 2 else AC2)

# ═══════════════════════════════════════════════════════════
# SOLUTION
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
multi_text(s, 80, 60, 800, [
    ("新架构：Grid Canvas", 32, DARK, True, PP_ALIGN.LEFT),
    ("Agent 选格子 · 引擎判冲突 · 通过才写 PPT", 18, GRAY, False, PP_ALIGN.LEFT),
])

# Two big cards for the two layers
card(s, 60, 160, 400, 300, "定位层 PositioningGrid",
     "16×9 粗网格, 60pt/cell, Excel 命名 (A1..P9)\n\n"
     "Agent 的空间词汇\n"
     "- \"body 放 A2:D5\" → 翻译为 (60,60,240,240) pt\n"
     "- cell_range: ['A1','B2'] → \"A1:B2\" 紧凑表达\n"
     "- 每 slide ~425 tokens",
     accent=AC1)
card(s, 500, 160, 400, 300, "信息层 InformationGrid",
     "32×18 细网格, 30pt/cell, 内部地图\n\n"
     "引擎的感知层\n"
     "- 每格带 owner_id, content_type, z_order\n"
     "- locked/source 标记模板装饰\n"
     "- checkpoint/rollback 操作前快照",
     accent=AC2)

# ═══════════════════════════════════════════════════════════
# INTERACTION MATRIX
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK)
rect(s, 0, 0, 960, 540, fill=DARK)
text_box(s, 80, 60, 800, 40, "交互矩阵 — 类型×类型 → 判定", size=28, color=RGBColor(0xFF,0xFF,0xFF), bold=True)
hline(s, 80, 110, 200)

matrix_data = [
    ("", "text", "textbox", "image", "table", "chart", "shape"),
    ("text", "❌", "✓", "❌", "❌", "❌", "✓"),
    ("textbox", "✓", "❌", "✓", "✓", "✓", "✓"),
    ("image", "❌", "✓", "✓", "❌", "❌", "✓"),
    ("table", "❌", "✓", "❌", "❌", "❌", "✓"),
    ("chart", "❌", "✓", "❌", "❌", "❌", "✓"),
    ("shape", "✓", "✓", "✓", "✓", "✓", "✓"),
]
cw, ch = 110, 42
ox, oy = 80, 150
for ri, row in enumerate(matrix_data):
    for ci, cell in enumerate(row):
        x, y = ox + ci*cw, oy + ri*ch
        if ri == 0 or ci == 0:
            text_box(s, x+8, y+10, cw-16, ch, cell, size=11, color=RGBColor(0xAA,0xBB,0xDD), bold=True, align=PP_ALIGN.CENTER)
        else:
            bg_c = RGBColor(0x20,0x40,0x30) if cell == "✓" else RGBColor(0x40,0x20,0x20)
            rect(s, x+4, y+2, cw-8, ch-4, fill=bg_c)
            text_box(s, x+8, y+10, cw-16, ch, cell, size=14,
                     color=RGBColor(0x4F,0xE0,0x80) if cell == "✓" else RGBColor(0xF0,0x60,0x60),
                     align=PP_ALIGN.CENTER)

text_box(s, 80, 480, 800, 30, "有背景色块(vs 已有元素)列。未列出组合 → DEFAULT_POLICY = ALLOW", size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════
# FLOW COMPARISON
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
text_box(s, 80, 60, 800, 40, "Agent 操作流对比", size=28, bold=True)
hline(s, 80, 110, 200)

# Before
rect(s, 40, 160, 420, 50, fill=RGBColor(0xFF,0xEE,0xEE))
text_box(s, 60, 170, 400, 30, "旧：事后补救", size=16, color=RED, bold=True)
flow_old = ["1. Agent 算 pt 坐标", "2. python-pptx 写入文件", "3. audit() 全量扫描", "4. 发现重叠 → 报 Issue", "5. Agent 重新算坐标", "6. 再写 → 再审计 → 可能新问题"]
for i, step in enumerate(flow_old):
    y = 230 + i*40
    rect(s, 60, y, 380, 32, fill=RGBColor(0xFC,0xF0,0xF0))
    text_box(s, 70, y+5, 360, 24, step, size=12, color=RGBColor(0xC0,0x40,0x40))

# After
rect(s, 500, 160, 420, 50, fill=RGBColor(0xEE,0xFF,0xEE))
text_box(s, 520, 170, 400, 30, "新：事前拦截", size=16, color=AC2, bold=True)
flow_new = ["1. Agent 选格子 \"A2:D5\"", "2. Grid.try_place() 判定", "3. ❌→原因+空闲建议 ✓→占用信息层", "4. Agent 调整 / 确认", "5. Grid.commit() 原子写PPT", "6. PPT 文件从未被污染"]
for i, step in enumerate(flow_new):
    y = 230 + i*40
    c = RGBColor(0xF0,0xFC,0xF0) if "❌" not in step else RGBColor(0xFF,0xF8,0xF0)
    text_c = AC2 if "✓" not in step else RGBColor(0x00,0xA0,0x70)
    rect(s, 520, y, 380, 32, fill=c)
    text_box(s, 530, y+5, 360, 24, step, size=12, color=text_c)

hline(s, 40, 500, 880)

# ═══════════════════════════════════════════════════════════
# CORE API
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK)
rect(s, 0, 0, 960, 540, fill=DARK)
text_box(s, 80, 60, 800, 40, "核心 API — 4 个原语", size=28, color=RGBColor(0xFF,0xFF,0xFF), bold=True)
hline(s, 80, 110, 200)

cards_api = [
    ("try_place(element_id, content_type, target_cells)",
     "→ PlacementResult {verdict, conflicts, z_hint, free_suggestion}\n\n"
     "Agent \"body 放 A2:D6\" → 引擎判 BLOCK/ALLOW\n"
     "BLOCK 时返回冲突原因 + 空闲区域建议",
     AC1),
    ("commit(ppt_path)",
     "→ {status: ok, path: ...}\n\n"
     "信息层→临时PPT→os.replace 原子替换\n"
     "成功才确认 checkpoint，失败 rollback",
     AC2),
    ("rollback()",
     "→ {status: rolled_back}\n\n"
     "信息层恢复到 checkpoint，PPT 不变\n"
     "Agent 误操作可安全撤销",
     RGBColor(0xE0,0x8F,0x00)),
    ("grid_snapshot(level)",
     "→ 分层输出 JSON (L0/L1/L2)\n\n"
     "L0: ~50t 幻灯片总览\n"
     "L1: ~100t 区域详情\n"
     "L2: ~60t 元素全貌",
     RGBColor(0x7B,0x4F,0xB8)),
]
for i, (title, body, accent) in enumerate(cards_api):
    x = 40 + (i % 2) * 450
    y = 160 + (i // 2) * 175
    card(s, x, y, 420, 155, title, body, accent=accent)

# ═══════════════════════════════════════════════════════════
# SUPPLY LAYERS
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
text_box(s, 80, 60, 800, 40, "Agent 输出 — 三层供给", size=28, bold=True)
hline(s, 80, 110, 200)

layers = [
    ("L0 总览 (~50 tokens)", "slide × 5", "zone: A1:H1(title) A2:D6(body) E2:H6(fig)\nfree: A8:D9, I1:N9\ndensity: 31.9%", AC1),
    ("L1 区域详情 (~100 tokens)", "按需展开", "zone: A2:D6, type: body\nneighbors: E2:H6(fig, 0pt gap)\nsub: A3:D4 conflict with fig", AC2),
    ("L2 元素全貌 (~60 tokens)", "精确定位", "id: s01, type: text, cells: A2:D3\nfine_cells: 32, font: 14pt\nconflicts: [s05(image, 25%)]", RGBColor(0x7B,0x4F,0xB8)),
]
for i, (title, scope, detail, accent) in enumerate(layers):
    y = 160 + i * 120
    rect(s, 80, y, 800, 100, fill=BGL)
    rect(s, 80, y, 6, 100, fill=accent)
    text_box(s, 110, y+12, 200, 30, title, size=16, bold=True)
    text_box(s, 330, y+12, 120, 24, scope, size=11, color=GRAY)
    text_box(s, 110, y+44, 740, 50, detail, size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════
# ENGINEERING REFERENCES
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK)
rect(s, 0, 0, 960, 540, fill=DARK)
text_box(s, 80, 60, 800, 40, "工程参考 — 4 系统 → 1 方案", size=28, color=RGBColor(0xFF,0xFF,0xFF), bold=True)
hline(s, 80, 110, 200)

refs = [
    ("Unity Physics2D", "Layer Collision Matrix", "→ 交互矩阵 (BLOCK_PAIRS)", AC1),
    ("CSS Grid Layout", "grid-template-areas 命名", "→ 定位层格子命名 (A1..P9)", AC2),
    ("PCB Design Rules", "online DRC + batch DRC", "→ try_place(实时) + audit(全局)", RGBColor(0xE0,0x8F,0x00)),
    ("Figma Auto Layout", "对齐 + 间距 + 约束", "→ 对齐建议 + 密度热力图", RGBColor(0x7B,0x4F,0xB8)),
]
for i, (src, src_detail, arrow, accent) in enumerate(refs):
    y = 160 + i * 90
    text_box(s, 80, y+10, 200, 30, src, size=18, color=accent, bold=True)
    text_box(s, 280, y+14, 250, 24, src_detail, size=13, color=GRAY)
    text_box(s, 540, y+10, 100, 30, "──▶", size=20, color=RGBColor(0xA0,0xA0,0xB0), align=PP_ALIGN.CENTER)
    text_box(s, 640, y+14, 250, 24, arrow, size=14, color=RGBColor(0xFF,0xFF,0xFF))

# ═══════════════════════════════════════════════════════════
# DEMO OUTPUT
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
text_box(s, 80, 60, 800, 40, "Live Demo — 实际运行输出", size=28, bold=True)
hline(s, 80, 110, 200)

demo_lines = [
    ("1. Put title on A1:B1...           → ALLOW", True),
    ("2. Put body on A2:D6...            → ALLOW", True),
    ("3. Put figure on E2:H6...          → ALLOW", True),
    ("4. TRY caption on body (C6:D6)     → BLOCK", False),
    ("   Conflict: text 叠 text", False),
    ("   Suggestion: D1:E1", False),
    ("5. Move caption to A8:C8...        → ALLOW", True),
    ("6. Commit: ok, 28639 bytes", True),
    ("", True),
    ("Grid state == PPT file. 零运行时不一致.", True),
]
y = 150
for line, ok in demo_lines:
    c = AC2 if ok else RED
    text_box(s, 100, y, 760, 26, line, size=14, color=c, font_name="Consolas")
    y += 30

# ═══════════════════════════════════════════════════════════
# TEST RESULTS
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s, DARK)
rect(s, 0, 0, 960, 540, fill=DARK)
text_box(s, 80, 60, 800, 40, "测试结果", size=28, color=RGBColor(0xFF,0xFF,0xFF), bold=True)
hline(s, 80, 110, 200)

test_data = [
    ("grid/ 新测试", "34 tests", "8+5+8+7+3+3 = 34 ✓", AC2),
    ("test_mcp.py", "19 tools", "19/19 all green", AC2),
    ("collab_test.py", "6 scenarios", "6/6 人机协同", AC2),
    ("validate.py", "detection", "100% 召回 / 83.3% 精确", AC2),
    ("零回归", "old engine", "deprecated, still passing", RGBColor(0xE0,0x8F,0x00)),
]
for i, (name, scope, result, accent) in enumerate(test_data):
    y = 160 + i * 72
    rect(s, 80, y, 800, 58, fill=RGBColor(0x24,0x24,0x3E))
    rect(s, 80, y, 6, 58, fill=accent)
    text_box(s, 110, y+6, 220, 26, name, size=16, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    text_box(s, 350, y+6, 200, 26, scope, size=14, color=GRAY)
    text_box(s, 560, y+6, 300, 26, result, size=14, color=accent)

# ═══════════════════════════════════════════════════════════
# END
# ═══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, 0, 0, 960, 540, fill=DARK)
rect(s, 0, 0, 960, 6, fill=AC1)
rect(s, 0, 534, 960, 6, fill=AC1)

multi_text(s, 80, 120, 800, [
    ("Agent 选格子 · 引擎判冲突 · 通过才写 PPT", 28, (0xFF,0xFF,0xFF), True, PP_ALIGN.CENTER),
])
bullets_end = [
    "定位层 16×9 — Agent 的空间词汇，~425 tokens / slide",
    "信息层 32×18 — 引擎的感知地图，每个格子有类型 + 锁定 + z-order",
    "交互矩阵 — 6 个 ContentType，BLOCK_PAIRS 只 8 对，其余全放行",
    "事前拦截 — try_place → commit，PPT 文件从未被污染",
    "分层供给 — L0(50t) → L1(100t) → L2(60t)，冲突自动聚合",
]
y = 200
for b in bullets_end:
    text_box(s, 120, y, 720, 30, f"▸ {b}", size=15, color=RGBColor(0xD0,0xD4,0xE0))
    y += 40

text_box(s, 80, 450, 800, 30, "ppt_reflex/grid/  |  python grid/examples/demo_complete.py", size=12,
         color=RGBColor(0x88,0x90,0xA0), align=PP_ALIGN.CENTER)
text_box(s, 80, 480, 800, 24, "GitHub: iOfficeAI/OfficeCLI inspired rendering · Unity Physics2D inspired collision matrix", size=10,
         color=RGBColor(0x66,0x70,0x80), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
prs.save(OUT)
print(f"PPT saved: {OUT}")
print(f"Slides: {len(prs.slides)}, Size: {os.path.getsize(OUT)} bytes")
