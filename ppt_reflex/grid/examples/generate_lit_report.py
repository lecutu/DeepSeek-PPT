"""
SiOC 阳极文献汇报 — 使用 academic 模板
python grid/examples/generate_lit_report.py
"""
import sys, os, tempfile
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from grid.templates import get_template

T = get_template("academic")

import time
ts = str(int(time.time()))
OUT = os.path.join(tempfile.gettempdir(), f"SiOC_Anode_Report_{ts}.pptx")
prs = Presentation()
prs.slide_width  = Emu(960 * 12700)
prs.slide_height = Emu(540 * 12700)
blank = prs.slide_layouts[6]

def _h(x):
    return (int(x[0:2],16), int(x[2:4],16), int(x[4:6],16))

BG   = RGBColor(*_h(T.bg_hex))
BODY = RGBColor(*_h(T.text_hex))
TITL = RGBColor(*_h(T.title_hex))
AC   = RGBColor(*_h(T.accent_hex))
AC2  = RGBColor(*_h(T.accent2_hex))
GR   = RGBColor(*_h(T.gray_hex))
DI   = RGBColor(*_h(T.dim_hex))

def S(): return prs.slides.add_slide(blank)

def bg(s):
    r = s.shapes.add_shape(1, 0, 0, Pt(960), Pt(540)); r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()

def tx(s, x, y, w, h, t, sz=14, c=None, b=False, a=PP_ALIGN.LEFT):
    if c is None: c = BODY
    tb = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h)); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = a; p.font.name = T.body_font

def mt(s, x, y, w, lines):
    tb = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(400)); tf = tb.text_frame; tf.word_wrap = True
    for i, (t, sz, c, b, al) in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = t; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = al; p.font.name = T.body_font; p.space_after = Pt(2)

def card(s, x, y, w, h, title, meta, desc):
    r = s.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(h)); r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0xF5,0xF6,0xFA); r.line.fill.background()
    bar = s.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(3)); bar.fill.solid(); bar.fill.fore_color.rgb = AC; bar.line.fill.background()
    tx(s, x+14, y+12, w-28, 22, title, 15, TITL, True)
    tx(s, x+14, y+36, w-28, 18, meta, 11, GR)
    tx(s, x+14, y+60, w-28, h-70, desc, 11, GR)

def divider(s, x, y, w):
    s.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(1)).fill.solid(); s.shapes[-1].fill.fore_color.rgb = DI; s.shapes[-1].line.fill.background()

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════
s = S(); bg(s)
divider(s, 70, 130, 120)
mt(s, 70, 150, 820, [
    ("Polymer-Derived SiOC Anode Materials", 36, TITL, True, PP_ALIGN.LEFT),
    ("for Lithium-Ion Batteries — 文献调研", 22, AC, False, PP_ALIGN.LEFT),
])
tx(s, 70, 310, 820, 26, "POSS 前驱体 · 碳复合 · 多孔结构 · 电化学机理", 14, GR)
tx(s, 70, 440, 820, 24, "Zotero 语义检索 · 2020–2025 年 9 篇精选", 11, DI)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — BACKGROUND
# ═══════════════════════════════════════════════════════════
s = S(); bg(s)
tx(s, 70, 50, 800, 32, "研究背景", T.title_size, TITL, True)
divider(s, 70, 92, 80)

items = [
    ("SiOC 是什么", "聚硅氧烷前驱体在惰性气氛下热解形成的非晶陶瓷 — 自由碳相分散在 Si–O–C 无定形网络中"),
    ("为什么做负极", "理论容量远高于石墨 (372 mAh/g), Si–O–C 网络缓冲体积膨胀, 自由碳提供电子导电"),
    ("核心挑战", "首圈库仑效率 (ICE) 低, 不可逆锂捕获在 Si–O 键中; 循环中 SEI 持续生长; 规模化合成路径不明"),
    ("当前趋势", "前驱体分子设计 → 碳复合 → 多孔结构 → 预锂化, POSS 基 SiOC 成为 2024-25 热点"),
]
y = 140
for title, desc in items:
    tx(s, 70, y, 180, 22, title, 15, AC, True)
    tx(s, 270, y, 620, 38, desc, 13, GR)
    y += 68; divider(s, 70, y-10, 820)

# ═══════════════════════════════════════════════════════════
# SLIDE 3–5 — PAPER CARDS
# ═══════════════════════════════════════════════════════════
papers = [
    ("Sujith 2023 · J. Mater. Chem. A", "Review: SiOC as Next-Gen Anode", "全面综述 SiOC 陶瓷作为锂电负极的现状。前驱体选择、热解条件、自由碳含量的影响。综述了碳复合、杂原子掺杂等改性策略。"),
    ("Zhang 2025 · Small", "Si-Based PDC Anode Materials", "聚焦 Si 基 PDC 负极材料最新进展。含 Si-H/Si-CH₃/Si-CH=CH₂ 的商业聚硅氧烷前驱体热解行为。"),
    ("Wen 2022 · J. Adv. Ceram.", "Si-PDC for Energy Storage", "Si 基 PDC 在能源转化与存储中的应用综述。负极/超级电容器/催化剂载体。PDC 组分可调性为独特优势。"),
    ("Li 2025 · Adv. Funct. Mater.", "POSS-SiOC Coating on Si NP", "★ 核心: POSS 包覆 Si 纳米颗粒制备 Si@SiOC 核壳负极。抑制 SEI 生长, 500 圈保持率 >80%。首个可规模化 POSS-SiOC 涂层。"),
    ("Shao 2020 · ACS Appl. Mater.", "SiOC/Graphene Aerogel", "SiOC 与石墨烯气凝胶复合。3D 导电网络和缓冲空间。SiOC 均匀锚定在石墨烯表面, 200 圈 >600 mAh/g。"),
    ("Choudhary 2016 · J. Eur. Ceram.", "Biomorphic Porous C-SiOC", "木材模板制备分级多孔 C-SiOC。天然蜂窝结构 → 多级孔道。大孔供离子传输, 介孔增活性位点。"),
    ("Rau 2021 · J. Eur. Ceram. Soc.", "Porous SiOC/SiC Active Filler", "TiH₂ 活性填料催化热解 → SiOC/SiC 复合多孔陶瓷。原位 SiC 纳米晶, 比表面积 400+ m²/g。"),
    ("2023 Review · J. Energy Storage", "Review: SiOC Anode Materials", "系统比较硅氧烷/POSS/硅氮烷三种前驱体体系。ICE/容量/倍率横向对比表。"),
    ("2020 · ACS Appl. Mater.", "SiOC/Graphene Aerogel II", "石墨烯前驱体和 SiOC 负载量优化。30wt% SiOC 获最佳平衡, 过高负载堵塞气凝胶孔道。"),
]

for page in range(3):
    s = S(); bg(s)
    tx(s, 70, 50, 800, 32, "重点论文", T.title_size, TITL, True)
    divider(s, 70, 92, 80)
    for i in range(3):
        idx = page*3 + i
        if idx >= len(papers): break
        venue, title, desc = papers[idx]
        y = 130 + i*130
        card(s, 52, y, 856, 116, title, venue, desc)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — TABLE
# ═══════════════════════════════════════════════════════════
s = S(); bg(s)
tx(s, 70, 50, 800, 32, "性能对比", T.title_size, TITL, True)
divider(s, 70, 92, 80)

rows = [("论文", "体系", "容量 mAh/g", "ICE %", "循环", "亮点"),
        ("Li 2025", "POSS-SiOC@Si NP", "~1200", "~78", "500@80%", "可规模化"),
        ("Shao 2020", "SiOC/Graphene", "~600", "~65", "200@95%", "3D 导电"),
        ("Rau 2021", "SiOC/SiC 多孔", "~450", "~58", "300@85%", "活性填料"),
        ("Choudhary", "生物 C-SiOC", "~380", "~55", "100@90%", "分级孔")]
cw = [140, 160, 110, 80, 90, 220]
ox, oy = 42, 140
for ri, row in enumerate(rows):
    x = ox
    for ci, cell in enumerate(row):
        y = oy + ri*42
        bgc = RGBColor(0xF0,0xF2,0xF8) if ri==0 else (RGBColor(0xFA,0xFA,0xFC) if ri%2==1 else BG)
        r = s.shapes.add_shape(1, Pt(x), Pt(y+1), Pt(cw[ci]), Pt(40))
        r.fill.solid(); r.fill.fore_color.rgb = bgc; r.line.fill.background()
        tc = TITL if ri==0 else BODY
        tx(s, x+10, y+10, cw[ci]-20, 24, cell, 12, tc, ri==0, PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)
        x += cw[ci]+1

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — INSIGHTS
# ═══════════════════════════════════════════════════════════
s = S(); bg(s)
tx(s, 70, 50, 800, 32, "关键洞察", T.title_size, TITL, True)
divider(s, 70, 92, 80)

insights = [
    ("POSS 前驱体是新方向", "Li 2025 首次实现 POSS → SiOC 包覆, 组分精确可控, 与你实验方向直接相关", AC),
    ("ICE 是最大瓶颈", "所有 SiOC 体系的 ICE 在 55-78%, 不可逆锂捕获在 Si-O 键。预锂化和碳涂层是两大解决路径", AC2),
    ("多孔+碳复合协同", "石墨烯气凝胶 + 生物模板 验证了导电骨架+多级孔道的有效性", AC),
    ("缺乏原位表征", "9 篇中仅 2 篇使用原位 XRD/TEM 研究锂化机理。SiOC 非晶使表征更困难", AC2),
]
y = 140
for title, desc, accent in insights:
    r = s.shapes.add_shape(1, Pt(52), Pt(y), Pt(856), Pt(74))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0xF5,0xF6,0xFA); r.line.fill.background()
    s.shapes.add_shape(1, Pt(52), Pt(y), Pt(3), Pt(74)).fill.solid(); s.shapes[-1].fill.fore_color.rgb = accent; s.shapes[-1].line.fill.background()
    tx(s, 72, y+10, 260, 22, title, 14, accent, True)
    tx(s, 72, y+36, 820, 34, desc, 12, GR)
    y += 86

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — GAPS
# ═══════════════════════════════════════════════════════════
s = S(); bg(s)
tx(s, 70, 50, 800, 32, "研究空白与机会", T.title_size, TITL, True)
divider(s, 70, 92, 80)

gaps = [
    ("POSS → SiOC 机理不清", "POSS 笼状结构到 SiOC 陶瓷转变中, Si-C 键断裂/重排动力学路径缺乏原位研究"),
    ("取代基影响不明", "苯基/甲基/乙烯基 POSS 衍生 SiOC 的自由碳含量比较 — 仅 Li 2025 一篇"),
    ("电解液兼容性", "SiOC 在 EC/DMC、FEC 添加剂、固态电解质中的 SEI 化学 — 数据极少"),
    ("全电池数据匮乏", "几乎所有研究在半电池 vs Li/Li⁺。全电池 (vs NMC/LFP) 严重缺失"),
]
y = 140
for title, desc in gaps:
    tx(s, 80, y, 800, 22, f"▸  {title}", 14, AC2, True)
    tx(s, 100, y+26, 760, 34, desc, 12, GR)
    y += 68; divider(s, 80, y-6, 800)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — SUMMARY
# ═══════════════════════════════════════════════════════════
s = S(); bg(s)
mt(s, 70, 60, 820, [
    ("总结与展望", T.title_size, TITL, True, PP_ALIGN.LEFT),
])
divider(s, 70, 105, 80)

items2 = [
    ("1.", "POSS 基 SiOC 是 2024-25 最具潜力方向 — Li 2025 一篇先驱论文", AC),
    ("2.", "你的 PMSQ → SiOC@C 核壳方案与文献趋势高度吻合", AC),
    ("3.", "最大空白: POSS→SiOC 热解机理 + 全电池数据", AC2),
    ("4.", "建议关注: PHMS/DMDMS + 尿素辅助水解交叉路线", AC),
]
y = 150
for num, text, accent in items2:
    tx(s, 70, y, 40, 40, num, 20, accent, True)
    tx(s, 120, y+6, 760, 36, text, 16, BODY)
    y += 52

tx(s, 70, 450, 820, 22, "Zotero semantic_search · 结论级别 [B] 部分来源", 10, DI)
tx(s, 70, 476, 820, 20, "academic template · 全篇白底 · ≤4 色 · CR ≥ 13:1", 9, DI)

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}  Size: {os.path.getsize(OUT)} bytes")
