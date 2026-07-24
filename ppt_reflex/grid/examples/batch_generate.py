"""
batch_generate.py — Generate 6 PPTs with different templates/scenarios.
Tests the full pipeline: grid canvas, collision matrix, templates, ImagePrompter.
"""
import sys, os, tempfile, time, json
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from grid.templates import get_template
from grid import GridCanvas, GridConfig, ContentType, Supply


def _h(x):
    return (int(x[0:2], 16), int(x[2:4], 16), int(x[4:6], 16))


RESULTS = {}


# ═══════════════════════════════════════════════════════════════
# PPT 1 — Business Report (business template)
# ═══════════════════════════════════════════════════════════════
def make_business_report():
    T = get_template("business")
    OUT = os.path.join(tempfile.gettempdir(), f"Business_Report_{int(time.time())}.pptx")
    prs = Presentation()
    prs.slide_width = Emu(960 * 12700)
    prs.slide_height = Emu(540 * 12700)
    blank = prs.slide_layouts[6]

    BG = RGBColor(*_h(T.bg_hex))
    TX = RGBColor(*_h(T.text_hex))
    TI = RGBColor(*_h(T.title_hex))
    AC = RGBColor(*_h(T.accent_hex))
    AC2 = RGBColor(*_h(T.accent2_hex))
    GR = RGBColor(*_h(T.gray_hex))

    def S():
        return prs.slides.add_slide(blank)

    def bg_fill(s):
        r = s.shapes.add_shape(1, 0, 0, Pt(960), Pt(540))
        r.fill.solid()
        r.fill.fore_color.rgb = BG
        r.line.fill.background()

    def tx(s, x, y, w, h, t, sz=14, c=None, b=False, a=PP_ALIGN.LEFT):
        if c is None:
            c = TX
        tb = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = b
        p.alignment = a
        p.font.name = T.body_font

    def card(s, x, y, w, h, title, value, accent=AC):
        r = s.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(h))
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFA)
        r.line.fill.background()
        tx(s, x + 16, y + 12, w - 32, 24, title, 12, GR)
        tx(s, x + 16, y + 40, w - 32, 36, value, 28, accent, True)

    def line(s, x, y, w):
        s.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(1)).fill.solid()
        s.shapes[-1].fill.fore_color.rgb = RGBColor(0xE0, 0xE4, 0xEB)
        s.shapes[-1].line.fill.background()

    # Slide 1 — Cover
    s = S(); bg_fill(s)
    line(s, 60, 200, 120)
    tx(s, 60, 230, 820, 60, "Q3 2026 Product Strategy", 38, TI, True)
    tx(s, 60, 300, 820, 30, "DeepSeek PPT Maker — AI-Powered Presentation Generation", 18, AC)
    tx(s, 60, 440, 820, 24, "Confidential  |  July 2026  |  Strategy Team", 11, GR)

    # Slide 2 — KPI Dashboard
    s = S(); bg_fill(s)
    tx(s, 60, 40, 820, 36, "KPI Dashboard", 26, TI, True)
    line(s, 60, 84, 80)
    kpis = [("Revenue", "$12.4M", AC), ("Growth", "+23.5%", AC2),
            ("NPS", "78", AC), ("Churn", "2.1%", RGBColor(0xE0, 0x4F, 0x5F))]
    for i, (k, v, c) in enumerate(kpis):
        card(s, 50 + i * 220, 120, 200, 100, k, v, c)

    # Slide 3 — Bullets
    s = S(); bg_fill(s)
    tx(s, 60, 40, 820, 36, "Strategic Priorities", 26, TI, True)
    line(s, 60, 84, 80)
    bullets = [
        "Expand to APAC markets — target 3 new countries by Q4",
        "Launch AI-powered analytics dashboard — beta in September",
        "Reduce infrastructure costs by 15% via cloud migration",
        "Grow enterprise customer base from 120 to 200 accounts",
        "Invest in developer ecosystem — SDK + API docs overhaul",
    ]
    y = 130
    for b in bullets:
        tx(s, 80, y, 820, 28, f"▸  {b}", 15, TX)
        y += 48

    # Slide 4 — Competition
    s = S(); bg_fill(s)
    tx(s, 60, 40, 820, 36, "Competitive Landscape", 26, TI, True)
    line(s, 60, 84, 80)
    headers = ["Competitor", "Strength", "Weakness", "Our Edge"]
    cw = [180, 240, 240, 200]
    for ci, h in enumerate(headers):
        x = 52 + sum(cw[:ci])
        tx(s, x + 10, 130, cw[ci] - 20, 30, h, 13, TI, True)
    data = [
        ["Company A", "Market share 35%", "Slow innovation cycle", "6-month faster release"],
        ["Company B", "Strong brand", "High pricing ($500/mo)", "3x cheaper, same quality"],
        ["Company C", "Best AI features", "No enterprise support", "24/7 dedicated support"],
        ["Company D", "Open source", "No hosted option", "Managed cloud + self-hosted"],
    ]
    for ri, row in enumerate(data):
        y = 172 + ri * 52
        bg_c = RGBColor(0xF5, 0xF7, 0xFC) if ri % 2 == 0 else BG
        r = s.shapes.add_shape(1, Pt(50), Pt(y), Pt(860), Pt(46))
        r.fill.solid(); r.fill.fore_color.rgb = bg_c; r.line.fill.background()
        for ci, cell in enumerate(row):
            x = 52 + sum(cw[:ci])
            tx(s, x + 10, y + 10, cw[ci] - 20, 26, cell, 12, TX)

    # Slide 5 — Closing
    s = S(); bg_fill(s)
    line(s, 60, 180, 120)
    tx(s, 60, 220, 820, 60, "Thank You", 40, TI, True)
    tx(s, 60, 300, 820, 30, "Questions & Discussion", 18, AC)
    tx(s, 60, 440, 820, 22, "DeepSeek PPT Maker  |  github.com/lecutu/deepseek-ppt-maker", 11, GR)

    prs.save(OUT)
    RESULTS["business_report"] = {"path": OUT, "slides": len(prs.slides), "size": os.path.getsize(OUT)}
    print(f"[OK] Business Report: {len(prs.slides)} slides, {os.path.getsize(OUT)} bytes")


# ═══════════════════════════════════════════════════════════════
# PPT 2 — Teaching Slides (teaching template)
# ═══════════════════════════════════════════════════════════════
def make_teaching_slides():
    T = get_template("teaching")
    OUT = os.path.join(tempfile.gettempdir(), f"Teaching_Materials_{int(time.time())}.pptx")
    prs = Presentation()
    prs.slide_width = Emu(960 * 12700)
    prs.slide_height = Emu(540 * 12700)
    blank = prs.slide_layouts[6]

    BG = RGBColor(*_h(T.bg_hex))
    TX = RGBColor(*_h(T.text_hex))
    TI = RGBColor(*_h(T.title_hex))
    AC = RGBColor(*_h(T.accent_hex))
    AC2 = RGBColor(*_h(T.accent2_hex))
    GR = RGBColor(*_h(T.gray_hex))

    def S():
        return prs.slides.add_slide(blank)

    def bg_fill(s):
        r = s.shapes.add_shape(1, 0, 0, Pt(960), Pt(540))
        r.fill.solid()
        r.fill.fore_color.rgb = BG
        r.line.fill.background()

    def tx(s, x, y, w, h, t, sz=14, c=None, b=False, a=PP_ALIGN.LEFT):
        if c is None:
            c = TX
        tb = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.bold = b; p.alignment = a; p.font.name = T.body_font

    def code_box(s, x, y, w, h, code_lines):
        r = s.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(h))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x28, 0x2C, 0x34)
        r.line.fill.background()
        for i, line in enumerate(code_lines):
            c = RGBColor(0x98, 0xC3, 0x79) if line.strip().startswith("#") else RGBColor(0xDC, 0xDF, 0xE4)
            tx(s, x + 16, y + 12 + i * 24, w - 32, 22, line, 11, c, False, PP_ALIGN.LEFT)

    # Slide 1 — Title
    s = S(); bg_fill(s)
    tx(s, 60, 160, 840, 72, "Introduction to Python", 40, TI, True)
    tx(s, 60, 250, 840, 30, "Lecture 3: Functions & Modules", 20, AC)
    tx(s, 60, 440, 840, 22, "CS 101  |  Fall 2026  |  Prof. Zhang", 12, GR)

    # Slide 2 — Learning Objectives
    s = S(); bg_fill(s)
    tx(s, 60, 40, 840, 36, "Learning Objectives", 26, TI, True)
    objs = [
        "Define and call functions with parameters and return values",
        "Understand variable scope (local vs global)",
        "Import and use Python standard library modules",
        "Create custom modules and packages",
        "Write docstrings and follow PEP 8 conventions",
    ]
    y = 120
    for obj in objs:
        tx(s, 80, y, 820, 28, f"●  {obj}", 16, TX)
        y += 52

    # Slide 3 — Function Syntax
    s = S(); bg_fill(s)
    tx(s, 60, 40, 840, 36, "Function Syntax", 26, TI, True)
    code_box(s, 60, 110, 840, 220, [
        "# Function definition",
        "def calculate_capacity(charge, mass):",
        '    """Calculate specific capacity in mAh/g."""',
        "    capacity = charge / mass",
        "    return capacity",
        "",
        "# Calling the function",
        "result = calculate_capacity(charge=180, mass=0.5)",
        "print(f\"Specific capacity: {result} mAh/g\")",
    ])
    tx(s, 60, 350, 840, 50, "Key points: def keyword, parameters (arguments), return statement, docstring", 14, GR)

    # Slide 4 — Import Modules
    s = S(); bg_fill(s)
    tx(s, 60, 40, 840, 36, "Importing Modules", 26, TI, True)
    code_box(s, 60, 110, 420, 260, [
        "# Standard library",
        "import json",
        "from datetime import datetime",
        "from pathlib import Path",
        "",
        "# Third-party",
        "import numpy as np",
        "import matplotlib.pyplot as plt",
    ])
    code_box(s, 500, 110, 400, 260, [
        "# Custom module",
        "from my_package.grid import (",
        "    GridCanvas,",
        "    GridConfig,",
        "    ContentType,",
        ")",
        "",
        "# Best practice:",
        "# Use explicit imports",
    ])

    # Slide 5 — Practice
    s = S(); bg_fill(s)
    tx(s, 60, 40, 840, 36, "In-Class Exercise", 26, TI, True)
    tx(s, 80, 120, 800, 80, 'Write a function "compute_statistics(data)" that takes a list of numbers and returns a dict with keys: mean, median, std_dev, min, max.', 16, TX)
    code_box(s, 60, 230, 840, 260, [
        "def compute_statistics(data):",
        '    """Return summary statistics for a list of numbers."""',
        "    from statistics import mean, median, stdev",
        "    return {",
        '        "mean": mean(data),',
        '        "median": median(data),',
        '        "std_dev": stdev(data),',
        '        "min": min(data),',
        '        "max": max(data),',
        "    }",
    ])
    tx(s, 80, 510, 800, 40, "Time: 10 minutes. Submit via GitHub Classroom.", 14, AC2)

    # Slide 6 — Summary
    s = S(); bg_fill(s)
    tx(s, 60, 120, 840, 60, "Key Takeaways", 34, TI, True)
    takeaways = [
        "Functions = reusable blocks of code",
        "Parameters are inputs, return is output",
        "Modules organize code into files",
        "import gives access to external functionality",
    ]
    y = 220
    for t in takeaways:
        tx(s, 100, y, 760, 36, f"✓  {t}", 18, TX)
        y += 56
    tx(s, 60, 460, 840, 24, "Next: Lecture 4 — Classes & Objects", 13, GR)

    prs.save(OUT)
    RESULTS["teaching"] = {"path": OUT, "slides": len(prs.slides), "size": os.path.getsize(OUT)}
    print(f"[OK] Teaching: {len(prs.slides)} slides, {os.path.getsize(OUT)} bytes")


# ═══════════════════════════════════════════════════════════════
# PPT 3 — Product Launch (product template, dark)
# ═══════════════════════════════════════════════════════════════
def make_product_launch():
    T = get_template("product")
    OUT = os.path.join(tempfile.gettempdir(), f"Product_Launch_{int(time.time())}.pptx")
    prs = Presentation()
    prs.slide_width = Emu(960 * 12700)
    prs.slide_height = Emu(540 * 12700)
    blank = prs.slide_layouts[6]

    BG = RGBColor(*_h(T.bg_hex))
    TX = RGBColor(*_h(T.text_hex))
    TI = RGBColor(*_h(T.title_hex))
    AC = RGBColor(*_h(T.accent_hex))
    AC2 = RGBColor(*_h(T.accent2_hex))
    GR = RGBColor(*_h(T.gray_hex))

    def S():
        return prs.slides.add_slide(blank)

    def bg_fill(s):
        r = s.shapes.add_shape(1, 0, 0, Pt(960), Pt(540))
        r.fill.solid()
        r.fill.fore_color.rgb = BG
        r.line.fill.background()

    def tx(s, x, y, w, h, t, sz=14, c=None, b=False, a=PP_ALIGN.LEFT):
        if c is None:
            c = TX
        tb = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.bold = b; p.alignment = a; p.font.name = T.body_font

    def feature_card(s, x, y, w, icon, title, desc):
        r = s.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(160))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)
        r.line.fill.background()
        tx(s, x + 20, y + 12, w - 40, 30, icon, 28, AC, True)
        tx(s, x + 20, y + 48, w - 40, 26, title, 18, TI, True)
        tx(s, x + 20, y + 82, w - 40, 66, desc, 13, GR)

    # Slide 1 — Hero
    s = S(); bg_fill(s)
    tx(s, 60, 120, 840, 72, "Nebula v2.0", 48, TI, True, PP_ALIGN.CENTER)
    tx(s, 60, 210, 840, 36, "The AI-Native Development Platform", 22, AC, False, PP_ALIGN.CENTER)
    tx(s, 60, 280, 840, 30, "Build · Ship · Scale  —  10x faster", 16, GR, False, PP_ALIGN.CENTER)
    tx(s, 60, 460, 840, 24, "Launch Event  |  July 2026  |  San Francisco", 11, GR, False, PP_ALIGN.CENTER)

    # Slide 2 — Problem / Solution
    s = S(); bg_fill(s)
    tx(s, 60, 40, 840, 36, "The Problem", 26, TI, True)
    problems = [
        "Developers spend 40% of time on boilerplate and configuration",
        "Context switching between 12+ tools per feature",
        "Onboarding new team members takes 3-4 weeks",
        "Documentation drifts from code within days",
    ]
    y = 100
    for p in problems:
        tx(s, 60, y, 380, 28, f"✗  {p}", 14, RGBColor(0xE0, 0x60, 0x60))
        y += 42

    tx(s, 500, 40, 400, 36, "Our Solution", 26, AC2, True)
    solutions = [
        "AI code generation with context-aware prompts",
        "Unified dashboard — one tool, zero context switch",
        "Instant dev environments — code in 30 seconds",
        "Auto-synced docs from code + AI review",
    ]
    y = 100
    for s_text in solutions:
        tx(s, 500, y, 400, 28, f"✓  {s_text}", 14, RGBColor(0x4F, 0xE0, 0x80))
        y += 42

    # Slide 3 — Features
    s = S(); bg_fill(s)
    tx(s, 60, 40, 840, 36, "Key Features", 26, TI, True)
    features = [
        ("🧠", "AI Copilot", "Real-time code suggestions, test generation, bug detection. Context-aware with full repo understanding."),
        ("⚡", "Instant Deploy", "One-click deployment to 50+ regions. Preview URLs for every branch. Automatic rollback."),
        ("📊", "Observability", "Distributed tracing, metrics, logs in one view. AI anomaly detection with <5 min MTTR."),
        ("🔒", "Zero-Trust Security", "Built-in secrets management, automatic dependency auditing, SOC 2 Type II certified."),
    ]
    for i, (icon, title, desc) in enumerate(features):
        x = 40 + (i % 2) * 450
        y = 100 + (i // 2) * 180
        feature_card(s, x, y, 420, icon, title, desc)

    # Slide 4 — Pricing
    s = S(); bg_fill(s)
    tx(s, 60, 40, 840, 36, "Pricing", 26, TI, True)
    plans = [
        ("Starter", "Free", "5 projects, 10GB storage, community support", AC2),
        ("Pro", "$29/mo", "Unlimited projects, 100GB, priority support", AC),
        ("Enterprise", "Custom", "SSO, audit logs, SLA, dedicated support", RGBColor(0xC0, 0x78, 0xFF)),
    ]
    for i, (name, price, desc, color) in enumerate(plans):
        x = 60 + i * 290
        r = s.shapes.add_shape(1, Pt(x), Pt(110), Pt(260), Pt(280))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E); r.line.fill.background()
        tx(s, x + 20, 130, 220, 30, name, 20, TI, True, PP_ALIGN.CENTER)
        tx(s, x + 20, 180, 220, 40, price, 28, color, True, PP_ALIGN.CENTER)
        tx(s, x + 20, 280, 220, 80, desc, 13, GR, False, PP_ALIGN.CENTER)

    # Slide 5 — CTA
    s = S(); bg_fill(s)
    tx(s, 60, 160, 840, 72, "Get Started Today", 42, TI, True, PP_ALIGN.CENTER)
    tx(s, 60, 250, 840, 36, "nebula.dev  →  Sign up free", 20, AC, False, PP_ALIGN.CENTER)
    tx(s, 60, 320, 840, 30, "100,000+ developers already on Nebula", 14, GR, False, PP_ALIGN.CENTER)
    tx(s, 60, 450, 840, 22, "DeepSeek PPT Maker  |  github.com/lecutu/deepseek-ppt-maker", 10, GR, False, PP_ALIGN.CENTER)

    prs.save(OUT)
    RESULTS["product_launch"] = {"path": OUT, "slides": len(prs.slides), "size": os.path.getsize(OUT)}
    print(f"[OK] Product Launch: {len(prs.slides)} slides, {os.path.getsize(OUT)} bytes")


# ═══════════════════════════════════════════════════════════════
# PPT 4 — Data Report (data_report template)
# ═══════════════════════════════════════════════════════════════
def make_data_report():
    T = get_template("data_report")
    OUT = os.path.join(tempfile.gettempdir(), f"Data_Report_{int(time.time())}.pptx")
    prs = Presentation()
    prs.slide_width = Emu(960 * 12700)
    prs.slide_height = Emu(540 * 12700)
    blank = prs.slide_layouts[6]

    BG = RGBColor(*_h(T.bg_hex))
    TX = RGBColor(*_h(T.text_hex))
    TI = RGBColor(*_h(T.title_hex))
    AC = RGBColor(*_h(T.accent_hex))
    AC2 = RGBColor(*_h(T.accent2_hex))
    GR = RGBColor(*_h(T.gray_hex))

    def S():
        return prs.slides.add_slide(blank)

    def bg_fill(s):
        r = s.shapes.add_shape(1, 0, 0, Pt(960), Pt(540))
        r.fill.solid()
        r.fill.fore_color.rgb = BG
        r.line.fill.background()

    def tx(s, x, y, w, h, t, sz=14, c=None, b=False, a=PP_ALIGN.LEFT):
        if c is None:
            c = TX
        tb = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.bold = b; p.alignment = a; p.font.name = T.body_font

    # Slide 1 — Cover
    s = S(); bg_fill(s)
    tx(s, 60, 180, 840, 60, "Monthly Analytics Report", 38, TI, True)
    tx(s, 60, 250, 840, 30, "July 2026  |  DeepSeek PPT Maker Metrics", 18, AC)
    tx(s, 60, 440, 840, 20, "Data period: Jun 1 – Jun 30, 2026", 11, GR)

    # Slide 2 — Summary Stats
    s = S(); bg_fill(s)
    tx(s, 60, 40, 820, 36, "Executive Summary", 26, TI, True)
    stats = [("Total Users", "12,847", "+18.3%"), ("Active Sessions", "4.2M", "+12.7%"),
             ("Avg. Time", "24.6 min", "+3.2%"), ("Bounce Rate", "23.1%", "-5.4%")]
    for i, (k, v, delta) in enumerate(stats):
        x = 50 + i * 220
        r = s.shapes.add_shape(1, Pt(x), Pt(120), Pt(200), Pt(100))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFA); r.line.fill.background()
        dc = AC2 if delta.startswith("+") else RGBColor(0x4F, 0xAE, 0x60)
        tx(s, x + 16, 132, 168, 22, k, 12, GR)
        tx(s, x + 16, 160, 168, 32, v, 24, TI, True)
        tx(s, x + 16, 198, 168, 18, delta, 11, dc, True)

    # Slide 3 — Table
    s = S(); bg_fill(s)
    tx(s, 60, 40, 820, 36, "Monthly Breakdown by Channel", 26, TI, True)
    headers = ["Channel", "Visitors", "% Change", "Revenue", "Conv. Rate"]
    cw = [180, 150, 130, 150, 150]
    for ci, h in enumerate(headers):
        x = 52 + sum(cw[:ci])
        tx(s, x + 10, 120, cw[ci] - 20, 28, h, 13, TI, True, PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
    rows = [
        ["Organic Search", "48,200", "+22%", "$124,500", "3.2%"],
        ["Paid Ads", "22,100", "+8%", "$89,200", "2.1%"],
        ["Social Media", "15,600", "+45%", "$52,800", "4.5%"],
        ["Email", "9,800", "-3%", "$31,400", "1.8%"],
        ["Direct", "18,300", "+12%", "$67,900", "2.7%"],
    ]
    for ri, row in enumerate(rows):
        y = 160 + ri * 48
        bg_c = RGBColor(0xF5, 0xF7, 0xFC) if ri % 2 == 0 else BG
        r = s.shapes.add_shape(1, Pt(50), Pt(y), Pt(860), Pt(42))
        r.fill.solid(); r.fill.fore_color.rgb = bg_c; r.line.fill.background()
        for ci, cell in enumerate(row):
            x = 52 + sum(cw[:ci])
            tx(s, x + 10, y + 9, cw[ci] - 20, 24, cell, 12, TX, False, PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # Slide 4 — Insights
    s = S(); bg_fill(s)
    tx(s, 60, 40, 820, 36, "Key Insights & Recommendations", 26, TI, True)
    insights = [
        ("Social Media ROI", "45% traffic growth vs +8% for paid. Shift 20% ad budget to organic + social.", AC),
        ("Mobile Conversion Gap", "Mobile users convert at 1.9% vs 3.8% desktop. Prioritize mobile checkout UX.", AC2),
        ("Email Re-engagement", "9,800 visitors but -3% trend. A/B test subject lines, segment by inactivity.", AC),
        ("Peak Hours", "70% of sessions occur 9AM-6PM weekdays. Schedule deployments for weekends.", AC2),
    ]
    y = 110
    for title, desc, accent in insights:
        tx(s, 70, y, 30, 30, "▸", 18, accent, True)
        tx(s, 105, y + 2, 220, 24, title, 15, accent, True)
        tx(s, 105, y + 30, 780, 22, desc, 13, GR)
        y += 72

    # Slide 5 — Next Steps
    s = S(); bg_fill(s)
    tx(s, 60, 40, 820, 36, "Next Steps", 26, TI, True)
    steps = [
        "1.  Launch mobile checkout redesign — target completion: Aug 15",
        "2.  Increase social media ad spend by 30% — test 3 new platforms",
        "3.  Implement email segmentation pipeline — engineering sprint next week",
        "4.  Schedule quarterly review with stakeholders — Aug 30",
    ]
    y = 120
    for step in steps:
        tx(s, 80, y, 800, 30, step, 16, TX)
        y += 52

    prs.save(OUT)
    RESULTS["data_report"] = {"path": OUT, "slides": len(prs.slides), "size": os.path.getsize(OUT)}
    print(f"[OK] Data Report: {len(prs.slides)} slides, {os.path.getsize(OUT)} bytes")


# ═══════════════════════════════════════════════════════════════
# PPT 5 — Collision pressure test (grid canvas)
# ═══════════════════════════════════════════════════════════════
def make_collision_test():
    """Deliberately stress collision matrix — all 6 content types, overlapping, blocking."""
    T = get_template("academic")
    OUT = os.path.join(tempfile.gettempdir(), f"Collision_Test_{int(time.time())}.pptx")
    prs = Presentation()
    prs.slide_width = Emu(960 * 12700)
    prs.slide_height = Emu(540 * 12700)
    blank = prs.slide_layouts[6]

    BG = RGBColor(*_h(T.bg_hex))
    TX = RGBColor(*_h(T.text_hex))
    TI = RGBColor(*_h(T.title_hex))
    AC = RGBColor(*_h(T.accent_hex))

    def S():
        return prs.slides.add_slide(blank)

    def bg_fill(s):
        r = s.shapes.add_shape(1, 0, 0, Pt(960), Pt(540))
        r.fill.solid()
        r.fill.fore_color.rgb = BG
        r.line.fill.background()

    def tx(s, x, y, w, h, t, sz=14, c=None, b=False):
        if c is None:
            c = TX
        tb = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
        p.font.name = T.body_font

    # Use GridCanvas for this one
    canvas = GridCanvas(GridConfig())
    supply = Supply(GridConfig())

    # Fill slide 0 with real content via grid + commit to PPT
    placements = [
        ("title", ContentType.TEXT, ["A1", "B1", "C1", "D1", "E1", "F1", "G1", "H1"]),
        ("subtitle", ContentType.TEXT, ["A2", "B2", "C2", "D2"]),
        ("body_left", ContentType.TEXT, ["A3", "B3", "C3", "D3", "A4", "B4", "C4", "D4",
                                          "A5", "B5", "C5", "D5", "A6", "B6", "C6", "D6"]),
        ("body_right", ContentType.TEXT, ["E3", "F3", "G3", "H3", "I3", "J3",
                                           "E4", "F4", "G4", "H4", "I4", "J4"]),
        ("figure_1", ContentType.IMAGE, ["A7", "B7", "C7", "D7", "A8", "B8", "C8", "D8"]),
        ("table_1", ContentType.TABLE, ["E7", "F7", "G7", "H7", "E8", "F8", "G8", "H8"]),
    ]

    accepted = 0
    rejected = 0
    for eid, ctype, cells in placements:
        r = canvas.try_place(eid, ctype, cells)
        if r.allowed:
            accepted += 1
        else:
            rejected += 1
            print(f"  [BLOCKED] {eid}: {supply.format_conflict(r).get('suggestions', ['?'])[0]}")

    # Now try deliberate collisions
    collisions_tried = [
        ("bad_overlap_text", ContentType.TEXT, ["C4", "D4"]),       # TEXT on TEXT → BLOCK
        ("bad_overlap_img", ContentType.IMAGE, ["A3", "B3"]),       # IMAGE on TEXT → BLOCK
        ("bad_overlap_table", ContentType.TABLE, ["A7", "B7"]),    # TABLE on IMAGE → BLOCK
        ("good_textbox", ContentType.TEXTBOX, ["C4", "D4"]),       # TEXTBOX on TEXT → ALLOW
    ]

    for eid, ctype, cells in collisions_tried:
        r = canvas.try_place(eid, ctype, cells)
        if r.allowed:
            accepted += 1
            print(f"  [ALLOW] {eid} ({ctype.name}) on existing — valid overlap")
        else:
            rejected += 1
            print(f"  [BLOCKED] {eid} ({ctype.name}) — expected: {r.conflicts[0].detail if r.conflicts else '?'}")

    canvas.checkpoint()
    canvas.commit(OUT)

    # Build slides manually for overview
    s = S(); bg_fill(s)
    tx(s, 60, 40, 840, 40, "Collision Matrix — Stress Test", 28, TI, True)
    tx(s, 60, 120, 840, 30, f"Accepted: {accepted}  |  Rejected: {rejected}", 18, AC, True)
    tx(s, 60, 180, 840, 30, "Grid canvas: try_place → commit workflow", 14, TX)

    grid_summary = [
        f"Grid density: {canvas.info_grid.density():.0%}",
        f"Matrix: {len(ContentType)} content types, 8 BLOCK_PAIRS",
        f"All accepted placements committed atomically — PPT file untouched on rollback.",
    ]
    y = 240
    for gs in grid_summary:
        tx(s, 80, y, 800, 26, f"▸  {gs}", 13, TX)
        y += 36

    prs.save(OUT)
    RESULTS["collision_test"] = {"path": OUT, "slides": len(prs.slides), "size": os.path.getsize(OUT)}
    print(f"[OK] Collision Test: {len(prs.slides)} slides, {os.path.getsize(OUT)} bytes, {accepted} accepted / {rejected} blocked")


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("=" * 60)
    print("Batch PPT Generation — 5 templates × different scenarios")
    print("=" * 60)

    make_business_report()
    make_teaching_slides()
    make_product_launch()
    make_data_report()
    make_collision_test()

    print()
    print("=" * 60)
    print("Results:")
    print("=" * 60)
    total_slides = 0
    for name, info in RESULTS.items():
        total_slides += info["slides"]
        print(f"  {name:20s}  {info['slides']} slides  {info['size']:>6} bytes")
    print(f"  {'TOTAL':20s}  {total_slides} slides")
    print(f"\n  All files in: {tempfile.gettempdir()}")
