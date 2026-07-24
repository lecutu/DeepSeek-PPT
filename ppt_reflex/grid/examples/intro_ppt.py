"""
examples/intro_ppt.py — 用 grid/ 制作新工具介绍 PPT

演示: 定位层→信息层→交互矩阵→try_place→commit 全流程
"""
import sys, os, tempfile
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = os.path.join(tempfile.gettempdir(), "PPT_Reflex_Grid_Intro.pptx")
prs = Presentation()
prs.slide_width = Emu(960 * 12700)
prs.slide_height = Emu(540 * 12700)
blank = prs.slide_layouts[6]

# ── utilities ─────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(blank)

def rect(slide, x, y, w, h, fill=None):
    s = slide.shapes.add_shape(1, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(*fill)
    s.line.fill.background()
    return s

def text(slide, x, y, w, h, txt, size=14, color=(255,255,255), bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*color)
    p.font.bold = bold
    p.alignment = align
    return tb

def title_slide(slide, title, subtitle=""):
    rect(slide, 0, 0, 960, 540, fill=(18, 18, 32))
    text(slide, 60, 180, 840, 80, title, size=36, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        text(slide, 60, 270, 840, 50, subtitle, size=18, color=(160,160,200), align=PP_ALIGN.CENTER)

def section_slide(slide, num, title):
    rect(slide, 0, 0, 960, 540, fill=(24, 24, 48))
    rect(slide, 60, 60, 840, 420, fill=(32, 32, 64))
    text(slide, 100, 100, 200, 40, f"§{num}", size=48, color=(100,180,255), bold=True)
    text(slide, 100, 200, 800, 60, title, size=28, bold=True)

def body_slide(slide, title, bullets):
    rect(slide, 0, 0, 960, 540, fill=(22, 22, 40))
    text(slide, 60, 40, 840, 50, title, size=24, bold=True, color=(100,180,255))
    y = 120
    for b in bullets:
        text(slide, 80, y, 800, 36, f"▸ {b}", size=16, color=(220,220,240))
        y += 50

def code_slide(slide, title, code_lines):
    rect(slide, 0, 0, 960, 540, fill=(16, 16, 32))
    text(slide, 60, 40, 840, 50, title, size=22, bold=True, color=(100,180,255))
    y = 110
    for line in code_lines:
        c = (180, 220, 140) if line.startswith("#") else (220, 220, 240)
        text(slide, 60, y, 840, 22, line, size=13, color=c)
        y += 26

def arch_slide(slide, title, boxes):
    rect(slide, 0, 0, 960, 540, fill=(22, 22, 40))
    text(slide, 60, 40, 840, 50, title, size=24, bold=True, color=(100,180,255))
    y = 130
    for label, desc, color in boxes:
        rect(slide, 100, y, 760, 55, fill=color)
        text(slide, 120, y+8, 300, 25, label, size=16, bold=True)
        text(slide, 420, y+8, 420, 25, desc, size=13, color=(200,200,220))
        y += 70

def comparison_slide(slide, title, rows):
    rect(slide, 0, 0, 960, 540, fill=(22, 22, 40))
    text(slide, 60, 40, 840, 50, title, size=24, bold=True, color=(100,180,255))
    # header
    text(slide, 60, 120, 280, 30, "之前", size=16, color=(255,100,100), bold=True, align=PP_ALIGN.CENTER)
    text(slide, 360, 120, 240, 30, "→", size=24, color=(200,200,200), align=PP_ALIGN.CENTER)
    text(slide, 620, 120, 280, 30, "现在", size=16, color=(100,255,100), bold=True, align=PP_ALIGN.CENTER)
    y = 170
    for before, after in rows:
        text(slide, 60, y, 280, 28, before, size=13, color=(220,180,180), align=PP_ALIGN.CENTER)
        text(slide, 620, y, 280, 28, after, size=13, color=(180,220,180), align=PP_ALIGN.CENTER)
        rect(slide, 60, y+30, 840, 1, fill=(60,60,80))
        y += 45


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — 封面
# ═══════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 960, 540, fill=(12, 12, 28))
rect(s, 0, 400, 960, 4, fill=(100, 180, 255))
text(s, 60, 160, 840, 100, "PPT Reflex Engine", size=44, bold=True, align=PP_ALIGN.CENTER)
text(s, 60, 270, 840, 50, "grid/ — 感知型网格画布", size=24, color=(100,180,255), align=PP_ALIGN.CENTER)
text(s, 60, 430, 840, 30, "两层网格 · 事前拦截 · 交互矩阵 · Agent 空间语言", size=16, color=(160,160,200), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — 问题
# ═══════════════════════════════════════════════════════════
body_slide(s, "为什么需要新架构？", [
    "python-pptx → 事后 audit → 发现重叠 → 报 issue → 重新改",
    "Agent 需自己算坐标、判断碰撞 → token 浪费、容易出 bugs",
    "PPT 文件已被污染 → 修复是事后弥补，不是事前预防",
    "跨 slide 持久化 bug: 修复在内存完成但无法可靠写入文件",
    "字体 / 间距 / 密度 — 引擎检测到但无法修复",
])

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — 两层架构
# ═══════════════════════════════════════════════════════════
s = add_slide()
section_slide(s, 1, "两层网格架构")

s2 = add_slide()
arch_slide(s2, "定位层 (Positioning Grid) — Agent 的空间词汇", [
    ("16×9 粗网格", "60pt/cell, Excel 命名 (A1..P9), ~425 tokens/slide", (30,50,80)),
    ("cells_to_bbox", "\"A2:D5\" → (60, 60, 240, 240) pt", (30,50,80)),
    ("bbox_to_cells", "(60,60,240,240) → [\"A2\",\"B2\",\"A3\",...]", (30,50,80)),
    ("cell_range", "紧凑表示: ['A1','A2','B1','B2'] → \"A1:B2\"", (30,50,80)),
])

s3 = add_slide()
arch_slide(s3, "信息层 (Information Grid) — 引擎的内部地图", [
    ("32×18 细网格", "30pt/cell, 引擎内部使用, Agent 不感知", (50,30,30)),
    ("InfoCell", "owner_id · content_type · z_order · locked · source", (50,30,30)),
    ("occupy/release", "占格 → 标记类型 → 后续 try_place 自动检测冲突", (50,30,30)),
    ("checkpoint/rollback", "操作前快照 → 失败回滚 → PPT 文件不变", (50,30,30)),
])

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — 交互矩阵
# ═══════════════════════════════════════════════════════════
s4 = add_slide()
section_slide(s4, 2, "交互矩阵 — 类型×类型 → 判定")

s5 = add_slide()
code_slide(s5, "BLOCK_PAIRS — 只拦真正出问题的", [
    "# 这 8 对类型重叠会被拦截:",
    "(TEXT, TEXT)       → BLOCK   # 文字叠文字 → 不可读",
    "(TEXT, IMAGE)      → BLOCK   # 文字叠图片 → 信息丢失",
    "(TEXT, TABLE)      → BLOCK",
    "(TEXT, CHART)      → BLOCK",
    "(IMAGE, TABLE)     → BLOCK",
    "(TABLE, TABLE)     → BLOCK   # 表格叠表格",
    "(CHART, CHART)     → BLOCK",
    "(TEXTBOX, TEXTBOX) → BLOCK   # 色块叠色块",
    "",
    "# 其余组合 → DEFAULT_POLICY = ALLOW",
    "# 文字叠色块 ✓ | 图片叠图片 ✓ | 标注叠任何 ✓",
])

s6 = add_slide()
code_slide(s6, "Z-Order 智能提示", [
    "# 引擎不只判 ALLOW/BLOCK，还告诉你谁在上:",
    "TEXT on TEXTBOX  → z_hint: new_above   # 文字浮在色块上",
    "TEXTBOX on TEXT  → z_hint: new_below   # 色块垫在文字下",
    "ANNOTATION on *  → z_hint: new_above   # 标注永远浮在最上",
    "IMAGE on TEXTBOX → z_hint: either      # 都可以",
    "",
    "# Agent 不需要自己判断 z-order",
    "# 引擎给出建议，Agent 选择策略",
])

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — 操作流程对比
# ═══════════════════════════════════════════════════════════
s7 = add_slide()
section_slide(s7, 3, "Agent 操作流 — 事前拦截 vs 事后补救")

s8 = add_slide()
comparison_slide(s8, "流程对比", [
    ("修改坐标 → 写 PPT → audit 扫描", "Agent 选格子 → Grid.try_place()"),
    ("发现重叠 → 报 issue → Agent 改", "❌ BLOCK → 原因+建议 → 调整"),
    ("再 audit → 可能有新问题", "✅ ALLOW → commit → PPT 文件"),
    ("PPT 可能已被污染", "PPT 从未被接触（原子写入）"),
    ("Agent 需理解 pt 坐标", "Agent 只需格子地址 + 内容类型"),
])

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — 核心 API
# ═══════════════════════════════════════════════════════════
s9 = add_slide()
section_slide(s9, 4, "核心 API — 4 个原语")

s10 = add_slide()
code_slide(s10, "GridCanvas.try_place()", [
    "result = canvas.try_place(",
    "    element_id=\"shape-07\",",
    "    content_type=ContentType.TEXT,",
    "    target_cells=[\"A2\",\"B2\",\"A3\",\"B3\"]",
    ")",
    "",
    "# result.verdict  → ALLOW | WARN | BLOCK",
    "# result.conflicts → [{cell, existing_id, type, verdict}, ...]",
    "# result.z_hint    → \"new_above\" | \"new_below\" | None",
    "# result.free_suggestion → [\"A8:D9\", \"E1:H2\"]",
])

s11 = add_slide()
code_slide(s11, "GridCanvas.commit() / rollback()", [
    "# 原子写入 — 临时文件 → os.replace → 确认",
    "canvas.checkpoint()",
    "canvas.commit(\"output.pptx\")",
    "# → Grid 状态 == PPT 文件，严格一致",
    "",
    "# 失败回滚 — PPT 不变",
    "canvas.rollback()",
    "# → 信息层恢复到 checkpoint 状态",
    "",
    "# 往返一致性:",
    "# grid → serializer.grid_to_ppt → PPT → serializer.ppt_to_grid → grid",
    "# 密度 0.134 ⇄ 0.134  ✓",
])

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — Agent 输出分层
# ═══════════════════════════════════════════════════════════
s12 = add_slide()
section_slide(s12, 5, "Agent 输出 — 三层供给 (Supply)")

s13 = add_slide()
code_slide(s13, "L0 总览 (~50 tokens) → L1 区域 → L2 元素", [
    "# L0: 幻灯片总览 — Agent 看到的第一眼",
    '{"slide": 5, "zones": {',
    '  "s_title": {"range": "A1:H1"},',
    '  "s_body":  {"range": "A2:D6"},',
    '  "s_fig":   {"range": "E2:H6"}',
    '}, "free": ["A8:D9", "I1:N9"], "density": 31.9}',
    "",
    "# L1: 区域详情 — 按需展开（~100 tokens）",
    "# L2: 元素全貌 — 单个元素的所有信息（~60 tokens）",
    "# 永远不全量导出 576 个信息格 — 冲突爆炸自动聚合",
])

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — 技术借鉴
# ═══════════════════════════════════════════════════════════
s14 = add_slide()
section_slide(s14, 6, "架构参考 — 4 个工程系统的精华")

s15 = add_slide()
body_slide(s15, "四参考系 → 一整合方案", [
    "Unity Physics2D Layer Collision Matrix → 交互矩阵 (BLOCK_PAIRS)",
    "CSS Grid Layout → 定位层命名区域 (A1:D5 = grid-area)",
    "PCB Design Rule Check → 双层检测: online DRC + batch DRC",
    "Figma Auto Layout → 对齐建议 + 间距分布",
])

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — Demo
# ═══════════════════════════════════════════════════════════
s16 = add_slide()
section_slide(s16, 7, "Live Demo — 5 步跑通")

s17 = add_slide()
code_slide(s17, "demo_3_elements.py — 输出截选", [
    "1. Put title on A1:B1...        → ALLOW",
    "2. Put body on A2:D6...         → ALLOW",
    "3. Put figure on E2:H6...       → ALLOW",
    "4. TRY caption on body (C6:D6)  → BLOCK",
    '   Conflicts: text 叠 text → 阻止',
    '   Suggestion: D1:E1',
    "5. Move caption to A8:C8...     → ALLOW",
    "6. Committed → 28639 bytes ✓",
    "",
    "# Grid state == PPT file. 零运行时不一致。",
])

# ═══════════════════════════════════════════════════════════
# SLIDE 10 — 测试结果
# ═══════════════════════════════════════════════════════════
s18 = add_slide()
section_slide(s18, 8, "测试 & 回归")

s19 = add_slide()
code_slide(s19, "34 tests + 25 regression = all green", [
    "# grid/ 新测试",
    "test_positioning.py     8/8   ✓   地址↔pt 转换",
    "test_matrix.py          5/5   ✓   交互矩阵判定",
    "test_canvas.py          8/8   ✓   try_place/commit/rollback",
    "test_serializer.py      7/7   ✓   Grid↔PPT 往返一致",
    "test_supply.py          3/3   ✓   输出 token 预算",
    "test_profiles.py        3/3   ✓   版式推断",
    "",
    "# 回归测试 (零破坏)",
    "test_mcp.py            19/19  ✓   全部 MCP 工具",
    "collab_test.py          6/6   ✓   人机协同 6 场景",
    "validate.py                100% 召回 / 83.3% 精确 ✓",
])

# ═══════════════════════════════════════════════════════════
# SLIDE 11 — 尾声
# ═══════════════════════════════════════════════════════════
s20 = add_slide()
rect(s20, 0, 0, 960, 540, fill=(12, 12, 28))
rect(s20, 0, 0, 960, 4, fill=(100, 180, 255))
rect(s20, 0, 536, 960, 4, fill=(100, 180, 255))
text(s20, 60, 130, 840, 60, "Agent 选格子 · 引擎判冲突 · 通过才写 PPT", size=28, bold=True, align=PP_ALIGN.CENTER)
text(s20, 60, 220, 840, 40, "ppt_reflex/grid/ — 感知型网格画布", size=20, color=(100,180,255), align=PP_ALIGN.CENTER)

bullets = [
    "定位层: 16×9 粗网格 — Agent 的空间词汇",
    "信息层: 32×18 细网格 — 引擎的感知层",
    "交互矩阵: 类型×类型 → BLOCK/ALLOW + z_hint",
    "事前拦截: try_place → commit, PPT 不被污染",
    "分层供给: L0(50t)/L1(100t)/L2(60t) Agent 输出",
]
y = 300
for b in bullets:
    text(s20, 120, y, 720, 30, f"▸ {b}", size=15, color=(200,200,220))
    y += 36

text(s20, 60, 500, 840, 25, "D:\\文献搜索员\\ppt_reflex\\grid\\  |  python grid/examples/demo_complete.py", size=12, color=(120,120,160), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
prs.save(OUT)
print(f"PPT saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(OUT)} bytes")
