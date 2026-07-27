"""
ppt_reflex/builder.py — Sole AI entry point. All engine capabilities exposed here, pure interface zero engine concepts.

from ppt_reflex.builder import PPTBuilder, load_style_presets, save_style_presets, list_style_presets

# Basic usage
builder = PPTBuilder(template="academic", style="academic_rigorous")
builder.add_slide("Cover",
    regions=[("r1", 100,80,760,380)],
    elements=[builder.title("Title"), builder.text("Body", style="Body")],
)
result = builder.build("out.pptx")

# Style preset management
presets = load_style_presets()           # read current presets
p = presets["academic_rigorous"]         # get one preset
p["color_override"]["bg"] = "#FFFDF5"    # change background
save_style_presets(presets)              # persist

# list_style_presets() -> [{id, display_name, mood, theme}, ...]
# lightweight guide for AI to pick a style without loading full presets
"""

from __future__ import annotations
import os, tempfile, time, math, json
from dataclasses import dataclass, field

from ppt_reflex.grid import (
    GridCanvas, GridConfig, ContentType, ElementPayload, Verdict,
    LayoutPlan, Region, Phase1Element, DecoIntent,
    execute_phase1, execute_phase2, global_composition_check,
)
from ppt_reflex.grid.templates import get_template, TemplateProfile
from ppt_reflex.grid.aesthetics import AestheticsEngine, AestheticViolation, ElemStyle
from ppt_reflex.grid.serializer import _render_image, _render_payload  # contain-fit rendering
from ppt_reflex.diff_log import DiffLog  # snapshot-based mutation trace, session lifetime
from ppt_reflex.roundtrip_check import check_overflow  # reopen saved PPTX and verify text fits
from ppt_reflex.color_triangulator import check_slide as tri_check_slide  # bg↔text↔fill color triangle

# ── Style presets path ──
_PRESETS_PATH = os.path.join(os.path.dirname(__file__), "style_presets.json")


def load_style_presets() -> dict:
    """Read style presets from disk. Returns full dict with meta + presets. Modify then call save_style_presets()."""
    with open(_PRESETS_PATH, "r", encoding="utf-8") as f:
        return json.load(f)


def save_style_presets(data: dict) -> None:
    """Save style presets to disk. Call load->modify->save."""
    with open(_PRESETS_PATH, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def list_style_presets() -> list[dict]:
    """Lightweight preset list for AI style selection. Each entry: id+display_name+mood+theme. No full color/shape data."""
    data = load_style_presets()
    return [
        {"id": pid, "display_name": p["display_name"], "mood": p["mood"], "theme": p["theme"]}
        for pid, p in data.get("presets", {}).items()
    ]


# ── WCAG luminance ──
def _lum(rgb: tuple) -> float:
    def f(c): s = c/255.0; return s/12.92 if s <= 0.04045 else ((s+0.055)/1.055)**2.4
    return 0.2126*f(rgb[0]) + 0.7152*f(rgb[1]) + 0.0722*f(rgb[2])

def _is_dark(rgb: tuple) -> bool:
    return _lum(rgb) < 0.25

def _hex_to_rgb(h: str) -> tuple:
    h = h.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _rgb_to_hex(rgb: tuple) -> str:
    """RGB tuple -> hex string (no # prefix), for AestheticsEngine."""
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

# ── Style table ──
STYLE = {
    "Heading":    dict(font_size=28, font_bold=True,  font_color=(0x1A,0x1A,0x2E), alignment="CENTER"),
    "Subtitle":   dict(font_size=18, font_color=(0x55,0x55,0x77), alignment="CENTER"),
    "Body":       dict(font_size=14, font_color=(0x33,0x33,0x44), font_name="Microsoft YaHei"),
    "Subheading": dict(font_size=16, font_bold=True,  font_color=(0x1B,0x3A,0x5C)),
    "Caption":    dict(font_size=10, font_color=(0x88,0x88,0x99)),
    "Footer":     dict(font_size=8,  font_color=(0xAA,0xAA,0xBB), alignment="CENTER"),
    "ListItem":   dict(font_size=13, font_color=(0x33,0x33,0x44), font_name="Microsoft YaHei"),
    "Emphasis":   dict(font_size=14, font_bold=True,  font_color=(0xC0,0x39,0x2B)),
}

# ── Built-in shape library ──
SHAPES = {
    "rounded_rectangle": "ROUNDED_RECTANGLE", "rectangle": "RECTANGLE",
    "oval": "OVAL", "parallelogram": "PARALLELOGRAM",
    "diamond": "DIAMOND", "chevron": "CHEVRON",
    "pentagon": "PENTAGON", "hexagon": "HEXAGON",
    "up_arrow": "UP_ARROW", "down_arrow": "DOWN_ARROW",
    "left_arrow": "LEFT_ARROW", "right_arrow": "RIGHT_ARROW",
    "star": "STAR_5_POINT", "triangle": "ISOSCELES_TRIANGLE",
    "home": "HOME_PLATE", "cross": "PLUS",
    "pie": "PIE", "wave": "WAVE", "donut": "DONUT",
    "plaque": "PLAQUE", "sun": "SUN",
}

# ── Internal spec ──
@dataclass
class _Spec:
    elem_id: str; style: str; text: str = ""; region: str = "main"
    ctype: str = "text"; fill_mode: str = "stack"
    pw: float|None = None; ph: float|None = None
    fill_color: tuple|None = None; shape_id: str = ""
    image_path: str = ""; margin: float = 6.0
    fit_mode: str = "fit"      # fit | fill | crop_center — fit=contain, no crop
    allow_upscale: bool = False # small images stay original size
    layout_mode: str = ""      # hero_top | hero_right | hero_left | center_float | small_inline | grid_2x2 | grid_1x3
    caption: str = ""          # Figure caption text
    # Phase1Element extended params (Fix #6)
    align_h: str = "left"
    allow_shrink: bool = False
    allow_wrap: bool = False
    arrow_slot: float = 48.0

@dataclass
class _Arrow:
    deco_id: str; from_elem: str; to_elem: str; text: str = ""
    direction: str = "below"; color: tuple = (0x66,0x66,0x66); width: float = 1.5
    # Fix #8: full DecoIntent params
    margin_pt: float = 8.0
    text_font_size: float = 10.0
    text_color: tuple = (0x55,0x55,0x55)
    occlusion_check: bool = True

@dataclass
class _Slide:
    title: str = ""
    regions: list = field(default_factory=list)
    elements: list[_Spec] = field(default_factory=list)
    arrows: list[_Arrow] = field(default_factory=list)


class PPTBuilder:
    """Sole AI entry point. add_slide -> build. Engine + templates fully transparent."""

    def __init__(self, template: str = "academic", style: str|None = None,
                 page_w: float = 960, page_h: float = 540,
                 template_pptx: str|None = None):
        self._t: TemplateProfile = get_template(template)
        self._style_preset: dict|None = None
        self._style_id: str|None = style
        self.pw, self.ph = page_w, page_h
        self._slides: list[_Slide] = []
        self._id = 0
        self._template_pptx = template_pptx
        self._style_body_font: str|None = None
        self._image_layout: dict|None = None
        self.diff_log = DiffLog()  # per-deck-session mutation trace, clear() on user confirm

        if style:
            data = load_style_presets()
            self._style_preset = data.get("presets", {}).get(style)
            if self._style_preset:
                c = self._style_preset["color_override"]
                fo = self._style_preset.get("font_override", {})
                so = self._style_preset.get("shape_override", {})
                overrides = dict(
                    bg_hex=c.get("bg", self._t.bg_hex),
                    text_hex=c.get("text_primary", self._t.text_hex),
                    title_hex=c.get("text_primary", self._t.title_hex),
                    accent_hex=c.get("accent", self._t.accent_hex),
                    accent2_hex=c.get("warn", self._t.accent2_hex),
                    gray_hex=c.get("text_secondary", self._t.gray_hex),
                    dim_hex=c.get("surface", self._t.dim_hex),
                    title_size=fo.get("scale_h1", self._t.title_size),
                    body_size=fo.get("scale_body", self._t.body_size),
                    caption_size=fo.get("scale_h2", self._t.caption_size),
                    divider_color_hex=c.get("accent", self._t.divider_color_hex),
                )
                self._t = self._t.override(**overrides)
                # Fix #14: capture body_font from style preset if available
                self._style_body_font = fo.get("body_font")
                # v2: capture image_layout from preset
                self._image_layout = self._style_preset.get("image_layout", None)

    def set_intent_scope(self, scope: dict):
        """Agent-declared intent scope: {"slide_ids":[2,3], "elem_ids":["box_3"]}.
        Used by DiffLog.scope_alert() to catch say-vs-do mismatch."""
        self.diff_log.set_intent_scope(scope)

    # ── slide ──
    def add_slide(self, title: str = "", *, regions: list|None = None,
                  elements: list|None = None, arrows: list|None = None) -> int:
        # Fix #5: regions can now have optional 6th field for content_inset
        if regions is None:
            regions = [("main", 60, 60, 840, 420, 1)]
        self._slides.append(_Slide(title, regions, elements or [], arrows or []))
        return len(self._slides) - 1

    def build_single_slide(self, slide_idx: int, prs=None) -> dict:
        """Build exactly one slide (for DeckPlanner harness). Returns per-slide diagnostics."""
        if slide_idx < 0 or slide_idx >= len(self._slides):
            return {"ok": False, "diagnostics": [{"kind": "invalid_index", "severity": "error",
                     "message": f"slide_idx {slide_idx} out of range [0, {len(self._slides)})"}]}
        spec = self._slides[slide_idx]
        plan = self._plan(spec)
        c = GridCanvas(GridConfig())
        c.checkpoint()

        self.diff_log.snap_before(plan, slide_idx)
        diags: list[dict] = []

        plan.validate(verbose=False)
        for d in plan.diagnostics:
            diags.append(_diag(slide_idx, "0.5", d))

        execute_phase1(plan, c)
        for d in plan.diagnostics:
            diags.append(_diag(slide_idx, "1", d))

        decos = execute_phase2(plan, c)
        for d in decos:
            if d.deco_type == "arrow" and d.x2:
                c.register_decoration(d.deco_id, "arrow", d.x1, d.y1, d.x2, d.y2,
                    line_color=d.style.get("line_color", (0x66,0x66,0x66)),
                    line_width_pt=d.style.get("line_width_pt", 1.5),
                    text=d.text, font_size=d.text_font_size, font_color=d.text_color)
            for w in d.occlusion_warnings:
                diags.append(_diag(slide_idx, "2", None, kind="arrow_occlusion",
                                   severity="warning", deco_id=d.deco_id, message=w))

        for ci in global_composition_check(plan):
            diags.append(_diag(slide_idx, "2.5", None, kind=ci.get("category","composition"),
                               severity=ci.get("level","info"), message=ci.get("message","")))

        ae_diags = self._run_aesthetics(c, plan)
        diags.extend(ae_diags)

        # P0-口②④: freeze → overflow check before render (single-slide path)
        from ppt_reflex.grid.text_metrics import check_overflow_2d
        for pe in plan.elements:
            p = pe.payload
            if not p or not p.text.strip():
                continue
            if pe.content_type not in (ContentType.TEXT, ContentType.TEXTBOX):
                continue
            v_auto_fit = not pe.height_is_locked
            h_auto_fit = not pe.width_is_locked
            issues = check_overflow_2d(
                p.text, p.font_size,
                box_w=pe.w, box_h=pe.h,
                line_spacing=p.line_spacing,
                v_auto_fit=v_auto_fit,
                h_auto_fit=h_auto_fit,
            )
            for iss in issues:
                severity = "warning" if pe.content_type == ContentType.TEXTBOX else iss["level"]
                diags.append({
                    "slide": slide_idx, "phase": "freeze",
                    "kind": iss["kind"],
                    "severity": severity,
                    "elem_id": pe.elem_id,
                    "message": iss["message"],
                })

        pv = c.pre_commit_validation()
        for err in pv.get("errors", []):
            diags.append(_diag(slide_idx, "pre", None, kind="validation_error", severity="error",
                               elem_id=err.get("owner_id",""), message=err.get("detail","")))
        for warn in pv.get("warnings", []):
            diags.append(_diag(slide_idx, "pre", None, kind="validation_warning", severity="warning",
                               elem_id=warn.get("owner_id",""), message=warn.get("detail","")))
        for adv in pv.get("advisories", []):
            diags.append(_diag(slide_idx, "pre", None, kind="advisory", severity="info",
                               elem_id=adv.get("owner_id",""), message=adv.get("detail","")))

        if prs is not None:
            _render_slide(prs, c, self._t, slide_index=slide_idx, total_slides=len(self._slides))
        self.diff_log.snap_after(plan, slide_idx)

        errs = [d for d in diags if d.get("severity") in ("error",)]
        warns = [d for d in diags if d.get("severity") in ("warning","warn")]
        return {"ok": len(errs) == 0, "diagnostics": diags,
                "summary": f"slide {slide_idx}: {len(diags)} issues ({len(errs)} errors, {len(warns)} warnings)",
                "roundtrip_ok": True}  # single-slide mode defers roundtrip to deck-level build()

    def build_stream(self, path: str|None = None):
        """分页流式构建 — 逐页 yield 诊断，AI 可以立即看到每页结果。

        用法：
            for slide_result in b.build_stream("out.pptx"):
                ok = slide_result["ok"]
                errs = [d for d in slide_result["diagnostics"] if d["severity"] == "error"]
                print(f"S{slide_result['slide']:02d}: {len(errs)} errors")
                if errs:
                    break  # 中断构建，只修这一页

        与 build() 的区别：
        - build() 一次性完成所有页 + roundtrip，返回 dict
        - build_stream() 逐页 yield，不阻塞等待全部完成
        - 第一个返回的是 {'type': 'start', ...} 元信息
        - 之后每个是 {'type': 'slide', ...} 单页诊断
        - 最后是 {'type': 'summary', ...} 总汇 (含 roundtrip)
        """
        from pptx import Presentation; from pptx.util import Pt
        if path is None:
            ts = int(time.time()); path = os.path.join(tempfile.gettempdir(), f"ppt_reflex_{ts}.pptx")

        if self._template_pptx and os.path.exists(self._template_pptx):
            prs = Presentation(self._template_pptx)
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        else:
            prs = Presentation()
        prs.slide_width = Pt(self.pw); prs.slide_height = Pt(self.ph)
        total_slides = len(self._slides)

        yield {"type": "start", "total_slides": total_slides, "template": self._t.id, "style": self._style_id}

        all_diags: list[dict] = []
        for i, spec in enumerate(self._slides):
            slide_result = self.build_single_slide(i, prs=prs)
            all_diags.extend(slide_result.get("diagnostics", []))
            yield {**slide_result, "type": "slide", "slide": i, "slide_total": total_slides}

        prs.save(path)
        errs = [d for d in all_diags if d.get("severity") in ("error",)]
        warns = [d for d in all_diags if d.get("severity") in ("warning","warn")]

        rt_results = check_overflow(path)
        for rt in rt_results:
            all_diags.append(_diag(rt["slide"], "rt", None,
                                   kind=rt["kind"], severity=rt["severity"],
                                   message=rt["message"]))
        rt_errors = [rt for rt in rt_results if rt["severity"] == "error"]
        rt_warns = [rt for rt in rt_results if rt["severity"] == "warning"]

        diff_report = self.diff_log.diff()
        scope = self.diff_log.scope_alert()
        yield {
            "type": "summary",
            "path": path, "ok": len(errs) == 0 and len(rt_errors) == 0,
            "diagnostics": all_diags,
            "summary": f"{len(all_diags)} issues ({len(errs)} errors, {len(warns)} warnings, "
                       f"{len(rt_errors)} roundtrip errors, {len(rt_warns)} roundtrip warnings)",
            "template": self._t.id, "style": self._style_id,
            "diff": {
                "entries": len(diff_report.entries) if diff_report else 0,
                "changed_elem_ids": list(diff_report.changed_elem_ids) if diff_report else [],
            },
            "scope_alert": scope,
        }

    def build(self, path: str|None = None) -> dict:
        from pptx import Presentation; from pptx.util import Pt
        if path is None:
            ts = int(time.time()); path = os.path.join(tempfile.gettempdir(), f"ppt_reflex_{ts}.pptx")

        # Use external template PPTX as base if provided (inherits master/layouts)
        if self._template_pptx and os.path.exists(self._template_pptx):
            prs = Presentation(self._template_pptx)
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        else:
            prs = Presentation()
        prs.slide_width = Pt(self.pw); prs.slide_height = Pt(self.ph)
        all_diags: list[dict] = []
        total_slides = len(self._slides)

        for i, spec in enumerate(self._slides):
            plan = self._plan(spec); c = GridCanvas(GridConfig()); c.checkpoint()

            # DiffLog: snap_before (entry — net change computed vs this baseline)
            self.diff_log.snap_before(plan, i)

            diags: list[dict] = []

            # Phase 0.5: region boundary validation
            plan.validate(verbose=False)
            for d in plan.diagnostics:
                diags.append(_diag(i, "0.5", d))

            # Phase 1: information layer layout
            execute_phase1(plan, c)
            for d in plan.diagnostics:
                diags.append(_diag(i, "1", d))

            # Phase 2: decoration layer resolution
            decos = execute_phase2(plan, c)
            for d in decos:
                if d.deco_type == "arrow" and d.x2:
                    c.register_decoration(d.deco_id, "arrow", d.x1, d.y1, d.x2, d.y2,
                        line_color=d.style.get("line_color", (0x66,0x66,0x66)),
                        line_width_pt=d.style.get("line_width_pt", 1.5),
                        text=d.text, font_size=d.text_font_size, font_color=d.text_color)
                for w in d.occlusion_warnings:
                    diags.append(_diag(i, "2", None, kind="arrow_occlusion", severity="warning",
                                       deco_id=d.deco_id, message=w))

            # Phase 2.5: global composition check
            for ci in global_composition_check(plan):
                diags.append(_diag(i, "2.5", None, kind=ci.get("category","composition"),
                                   severity=ci.get("level","info"), message=ci.get("message","")))

            ae_diags = self._run_aesthetics(c, plan)
            diags.extend(ae_diags)

            # Color triangle: bg ↔ text ↔ fill constraint system (Phase 3.0)
            tri_elems = []
            for pe in plan.elements:
                p = pe.payload
                if not p or not p.text.strip():
                    continue
                fc = p.font_color if p else (0xFF, 0xFF, 0xFF)
                fill = p.fill_color  # None → transparent
                tri_elems.append({
                    "elem_id": pe.elem_id,
                    "font_size": p.font_size if p else 14,
                    "font_bold": p.font_bold if p else False,
                    "font_color_rgb": fc,
                    "fill_color_rgb": fill,
                })
            tri_issues = tri_check_slide(tri_elems, self._t)
            for ti in tri_issues:
                diags.append({
                    "slide": i, "phase": "3.0", "kind": f"tri_{ti.edge.replace('↔','_')}",
                    "severity": ti.level, "elem_id": ti.elem_id,
                    "message": ti.message,
                })

            # P0-口②④: freeze → overflow check before render
            # Phase 1 locked coordinates (height_is_locked=True for stack/inline elements).
            # Aesthetics resolved final font_size. Freeze snapshots both, then 2D overflow
            # check reads the frozen state — same data renderer will use.
            # SHAPE_TO_FIT_TEXT cannot rescue overflow when layout has fixed the box dims.
            from ppt_reflex.grid.text_metrics import check_overflow_2d
            for pe in plan.elements:
                p = pe.payload
                if not p or not p.text.strip():
                    continue
                if pe.content_type not in (ContentType.TEXT, ContentType.TEXTBOX):
                    continue
                # v_auto_fit: SHAPE_TO_FIT_TEXT only actually grows if layout didn't lock height
                v_auto_fit = not pe.height_is_locked
                h_auto_fit = not pe.width_is_locked
                issues = check_overflow_2d(
                    p.text, p.font_size,
                    box_w=pe.w, box_h=pe.h,
                    line_spacing=p.line_spacing,
                    v_auto_fit=v_auto_fit,
                    h_auto_fit=h_auto_fit,
                )
                for iss in issues:
                    severity = "warning" if pe.content_type == ContentType.TEXTBOX else iss["level"]
                    diags.append({
                        "slide": i, "phase": "freeze",
                        "kind": iss["kind"],
                        "severity": severity,
                        "elem_id": pe.elem_id,
                        "message": iss["message"],
                    })

            # Fix #9: pre_commit_validation — bounds/overflow/role conflicts
            pv = c.pre_commit_validation()
            for err in pv.get("errors", []):
                diags.append(_diag(i, "pre", None, kind="validation_error", severity="error",
                                   elem_id=err.get("owner_id",""), message=err.get("detail","")))
            for warn in pv.get("warnings", []):
                diags.append(_diag(i, "pre", None, kind="validation_warning", severity="warning",
                                   elem_id=warn.get("owner_id",""), message=warn.get("detail","")))
            for adv in pv.get("advisories", []):
                diags.append(_diag(i, "pre", None, kind="advisory", severity="info",
                                   elem_id=adv.get("owner_id",""), message=adv.get("detail","")))

            # Fix #2: Render with smart layout selection
            _render_slide(prs, c, self._t, slide_index=i, total_slides=total_slides)

            # DiffLog: snap_after (exit — diff sees only net change, engine loop noise folded away)
            self.diff_log.snap_after(plan, i)

            all_diags.extend(diags)

        prs.save(path)
        errs = [d for d in all_diags if d.get("severity") in ("error",)]
        warns = [d for d in all_diags if d.get("severity") in ("warning","warn")]

        # Post-render roundtrip: reopen saved PPTX and verify every text box fits
        rt_results = check_overflow(path)
        for rt in rt_results:
            all_diags.append(_diag(rt["slide"], "rt", None,
                                   kind=rt["kind"], severity=rt["severity"],
                                   message=rt["message"]))
        rt_errors = [rt for rt in rt_results if rt["severity"] == "error"]
        rt_warns = [rt for rt in rt_results if rt["severity"] == "warning"]

        diff_report = self.diff_log.diff()
        scope = self.diff_log.scope_alert()
        return {"path": path, "ok": len(errs) == 0 and len(rt_errors) == 0,
                "diagnostics": all_diags,
                "summary": f"{len(all_diags)} issues ({len(errs)} errors, {len(warns)} warnings, "
                           f"{len(rt_errors)} roundtrip errors, {len(rt_warns)} roundtrip warnings)",
                "template": self._t.id, "style": self._style_id,
                "diff": {
                    "entries": len(diff_report.entries) if diff_report else 0,
                    "changed_elem_ids": list(diff_report.changed_elem_ids) if diff_report else [],
                },
                "scope_alert": scope}

    def clear_diff_log(self):
        """User confirmed deck is done — wipe mutation trace."""
        self.diff_log.clear()
    def title(self, text: str, region: str = "main") -> _Spec:
        return self._s("Heading", text, region, "text", ph=40)
    def subtitle(self, text: str, region: str = "main") -> _Spec:
        return self._s("Subtitle", text, region, "text", ph=30)
    def text(self, text: str, style: str = "Body", region: str = "main") -> _Spec:
        return self._s(style, text, region, "text")
    def bullet(self, text: str, region: str = "main") -> _Spec:
        return self._s("ListItem", f"• {text}", region, "text")
    def footer(self, text: str, region: str = "footer") -> _Spec:
        return self._s("Footer", text, region, "footer")
    def box(self, text: str, style: str = "Body", region: str = "main",
            fill_color: tuple|None = None, shape_id: str = "rounded_rectangle",
            ph: float|None = None, align_h: str = "left", allow_shrink: bool = False) -> _Spec:
        return self._s(style, text, region, "textbox", fill_color=fill_color,
                       shape_id=shape_id, ph=ph, align_h=align_h, allow_shrink=allow_shrink)
    def shape(self, shape_id: str, region: str = "main",
              fill_color: tuple|None = None, pw: float|None = None, ph: float|None = None) -> _Spec:
        return _Spec(elem_id=self._nid("shape"), style="", text="", region=region,
                     ctype="shape", fill_color=fill_color, shape_id=shape_id, pw=pw, ph=ph)
    def image(self, path: str, region: str = "main",
              pw: float|None = None, ph: float|None = None,
              fit_mode: str = "fit", allow_upscale: bool = False,
              layout_mode: str = "", caption: str = "") -> _Spec:
        # Fix #10: validate path exists
        if not os.path.isfile(path):
            print(f"[PPTBuilder] WARNING: image path not found: {path}")
        return _Spec(elem_id=self._nid("img"), style="", text="", region=region,
                     ctype="image", pw=pw, ph=ph, image_path=path,
                     fit_mode=fit_mode, allow_upscale=allow_upscale,
                     layout_mode=layout_mode, caption=caption)
    def arrow(self, frm: str, to: str, text: str = "", direction: str = "below",
              color: tuple = (0x66,0x66,0x66), width: float = 1.5,
              margin_pt: float = 8.0, text_font_size: float = 10.0,
              text_color: tuple = (0x55,0x55,0x55),
              occlusion_check: bool = True) -> _Arrow:
        # Allow passing _Spec objects directly — resolve to elem_id
        from_eid = frm.elem_id if hasattr(frm, 'elem_id') else frm
        to_eid = to.elem_id if hasattr(to, 'elem_id') else to
        return _Arrow(self._nid("arrow"), from_eid, to_eid, text, direction, color, width,
                      margin_pt, text_font_size, text_color, occlusion_check)
    def divider(self, region: str = "main", color: tuple|None = None, width_pt: float = 3.0) -> _Spec:
        dh = self._t.divider_color_hex
        if dh:
            c = color or tuple(int(dh.lstrip("#")[i:i+2], 16) for i in (0, 2, 4))
        else:
            c = color or (0x44, 0x44, 0x55)
        return _Spec(elem_id=self._nid("div"), style="", text="", region=region,
                     ctype="shape", fill_color=c, shape_id="rectangle", ph=width_pt)

    # ── Image layout auto-inference ──
    def auto_layout_mode(self, image_path: str) -> str:
        """Select layout mode from image aspect ratio + preset constraints. Decision tree:
        aspect > 1.6 -> hero_top (landscape)
        aspect < 0.8 -> hero_right (portrait)
        aspect 0.8-1.6 -> center_float (square)
        Falls back to preset's preferred_modes[0] if decision mode not in allowed set."""
        from PIL import Image
        try:
            img = Image.open(image_path)
            w, h = img.size
            aspect = w / h if h > 0 else 1.0
        except Exception as e:
            import traceback
            traceback.print_exc()
            print(f"[PPTBuilder] auto_layout_mode PIL error for {image_path}: {e}")
            return "center_float"

        # Decide by aspect ratio
        if aspect > 1.6:
            mode = "hero_top"
        elif aspect < 0.8:
            mode = "hero_right"
        else:
            mode = "center_float"

        # If preset constrains preferred_modes and decision mode is not in it, fall back to preferred_modes[0]
        if self._image_layout:
            preferred = self._image_layout.get("preferred_modes", [])
            if preferred and mode not in preferred:
                mode = preferred[0]
        return mode

    def image_constraints(self, layout_mode: str) -> dict:
        """Return max_w/max_h/anchor/ratio constraints for the given layout_mode from current preset."""
        mc = self._image_layout.get("mode_constraints", {}) if self._image_layout else {}
        return mc.get(layout_mode, {})

    def image_treatment(self) -> dict:
        """Return image rendering treatment (corner_radius/border/shadow) from current preset."""
        defaults = {"corner_radius_pt": 0, "border_role": "none", "shadow_role": "none"}
        if self._image_layout:
            defaults.update(self._image_layout.get("treatment", {}))
        return defaults

    def caption_format(self) -> dict:
        """Return caption format convention from current preset."""
        defaults = {"font_size": 11, "alignment": "left", "max_lines": 1, "prefix": ""}
        if self._image_layout:
            defaults.update(self._image_layout.get("caption", {}))
        return defaults

    # ── Internals ──
    def _nid(self, pfx: str) -> str:
        self._id += 1; return f"{pfx}_{self._id}"

    def _s(self, style: str, text: str, region: str, ctype: str,
           fill_color=None, shape_id="", ph=None,
           align_h="left", allow_shrink=False, allow_wrap=False) -> _Spec:
        s = STYLE.get(style, STYLE["Body"])
        return _Spec(self._nid("e"), style, text, region, ctype, "stack",
                     fill_color=fill_color or s.get("fill_color"),
                     shape_id=shape_id, ph=ph, margin=4.0,
                     align_h=align_h, allow_shrink=allow_shrink, allow_wrap=allow_wrap)

    # Fix #1: resolve style colors + font sizes from template profile
    def _resolve_style(self, style_name: str, fill_color=None) -> dict:
        """Merge STYLE defaults with template/style-preset colors AND font sizes."""
        s = dict(STYLE.get(style_name, STYLE["Body"]))
        t = self._t

        # Map style -> template color field
        if style_name in ("Heading", "Subheading"):
            s["font_color"] = _hex_to_rgb(t.title_hex)
        elif style_name == "Subtitle":
            s["font_color"] = _hex_to_rgb(t.gray_hex) if t.gray_hex else _hex_to_rgb(t.text_hex)
        elif style_name == "Emphasis":
            s["font_color"] = _hex_to_rgb(t.accent2_hex) if t.accent2_hex else _hex_to_rgb(t.accent_hex)
        elif style_name == "Caption":
            s["font_color"] = _hex_to_rgb(t.gray_hex) if t.gray_hex else _hex_to_rgb(t.text_hex)
        elif style_name == "Footer":
            s["font_color"] = _hex_to_rgb(t.dim_hex) if t.dim_hex else _hex_to_rgb(t.gray_hex)
        else:  # "Body", "ListItem"
            s["font_color"] = _hex_to_rgb(t.text_hex)

        # Map style -> template font size (was dead code — STYLE dict always won)
        if style_name in ("Heading",):
            s["font_size"] = t.title_size
        elif style_name in ("Subtitle",):
            s["font_size"] = t.caption_size
        elif style_name in ("Body", "ListItem", "Emphasis"):
            s["font_size"] = t.body_size
        elif style_name == "Caption":
            s["font_size"] = t.caption_size
        elif style_name == "Footer":
            s["font_size"] = t.page_number_size

        # Fix #13: font_name fallback
        if "font_name" not in s or not s["font_name"]:
            s["font_name"] = self._style_body_font or t.body_font

        # Dark fill -> white text
        if fill_color and _is_dark(fill_color):
            s["font_color"] = (0xFF, 0xFF, 0xFF)

        return s

    def _plan(self, spec: _Slide) -> LayoutPlan:
        ctmap = {"text": ContentType.TEXT, "textbox": ContentType.TEXTBOX,
                 "shape": ContentType.SHAPE, "image": ContentType.IMAGE,
                 "annotation": ContentType.ANNOTATION, "footer": ContentType.FOOTER}
        regions = []
        for ri, d in enumerate(spec.regions):
            rid, x, y, w, h = d[0], d[1], d[2], d[3], d[4]
            ro = d[5] if len(d) > 5 else ri + 1
            # Fix #5: optional 6th field for content_inset
            inset = d[6] if len(d) > 6 else 8.0
            regions.append(Region(rid, x, y, w, h, rid, ro, inset))

        elems = []
        for e in spec.elements:
            # Fix #1: resolve style from template colors
            s = self._resolve_style(e.style, e.fill_color)
            ctyp = ctmap.get(e.ctype, ContentType.TEXT)

            # Fix #4: set ElementPayload.role based on content_type
            role = None
            if e.ctype == "text":
                role = None  # handrail: TEXT -> ENTITY
            elif e.ctype == "textbox":
                role = None  # handrail: BAND -> ENTITY
            elif e.ctype == "shape":
                role = None  # handrail: BAND -> ENTITY (shape like divider)

            p = ElementPayload(
                role=role,
                text=e.text,
                font_size=s.get("font_size", 14),
                font_color=s.get("font_color", (0x33, 0x33, 0x44)),
                font_bold=s.get("font_bold", False),
                font_name=s.get("font_name", self._t.body_font),
                alignment=s.get("alignment", "LEFT"),
                fill_color=e.fill_color or s.get("fill_color"),
                shape_id=e.shape_id,
                line_spacing=self._t.line_spacing,
                image_path=e.image_path,
                fit_mode=e.fit_mode,
                allow_upscale=e.allow_upscale,
                layout_mode=e.layout_mode,
                caption=e.caption,
            )
            # Fix #6: expose Phase1Element params
            elems.append(Phase1Element(
                e.elem_id, e.region, ctyp,
                payload=p, fill_mode=e.fill_mode, margin_above=e.margin,
                preferred_width=e.pw, preferred_height=e.ph,
                align_h=e.align_h,
                allow_shrink=e.allow_shrink,
                allow_wrap=e.allow_wrap,
                ARROW_SLOT=e.arrow_slot,
            ))

        decos = []
        for a in spec.arrows:
            # Fix #8: pass all DecoIntent params
            decos.append(DecoIntent(
                a.deco_id, "arrow", [a.from_elem, a.to_elem],
                a.direction,
                margin_pt=a.margin_pt,
                style={"line_color": a.color, "line_width_pt": a.width},
                text=a.text,
                text_font_size=a.text_font_size,
                text_color=a.text_color,
                occlusion_check=a.occlusion_check,
            ))

        # Fix #7: use template page_margin for page_safe_inset
        plan = LayoutPlan(
            page_w=self.pw, page_h=self.ph,
            page_safe_inset=max(12.0, self._t.page_margin),
            title=spec.title,
            regions=regions,
            phase1_elements=elems,
            deco_intents=decos,
        )
        return plan

    def _run_aesthetics(self, canvas, plan) -> list[dict]:
        """Run AestheticsEngine and return structured diagnostics."""
        engine = AestheticsEngine()
        elems = []
        for pe in plan.elements:
            p = pe.payload
            fill_hex = _rgb_to_hex(p.fill_color) if p and p.fill_color else self._t.bg_hex.lstrip("#")
            font_hex = _rgb_to_hex(p.font_color) if p and p.font_color else "000000"
            es = ElemStyle(
                id=pe.elem_id, content_type=pe.content_type,
                font_size=p.font_size if p else 12,
                font_bold=p.font_bold if p else False,
                font_color=font_hex, fill_color=fill_hex,
                line_spacing=p.line_spacing if p else 1.2,
                text=p.text if p else "",
                x=pe.x, y=pe.y, w=pe.w, h=pe.h,
                auto_size="SHAPE_TO_FIT_TEXT",
                canvas_w=self.pw, canvas_h=self.ph,
            )
            # No fill -> inherit slide background (avoids false contrast errors when text over slide bg)
            if p and not p.fill_color and not p.shape_id:
                es.fill_color = self._t.bg_hex
            if p and p.fill_color:
                es.fill_color = _rgb_to_hex(p.fill_color)
            elems.append(es)

        violations = engine.check(elems, timing="audit")
        return [_ae_violation_to_diag(v) for v in violations]


def _diag(slide_idx, phase, d, kind="", severity="", deco_id="", elem_id="", message=""):
    """Normalize diagnostics from various sources into a uniform format."""
    if isinstance(d, dict):
        return {"slide": slide_idx, "phase": phase, **d}
    if hasattr(d, 'kind'):
        return {
            "slide": slide_idx, "phase": phase,
            "kind": d.kind, "severity": d.severity,
            "region_id": getattr(d, 'region_id', ''),
            "elem_id": getattr(d, 'elem_id', ''),
            "demand_pt": getattr(d, 'demand_pt', 0),
            "usable_pt": getattr(d, 'usable_pt', 0),
            "message": getattr(d, 'message', ''),
            "options": getattr(d, 'options', []),
        }
    if kind:
        return {"slide": slide_idx, "phase": phase, "kind": kind, "severity": severity,
                "deco_id": deco_id, "elem_id": elem_id, "message": message}
    return {"slide": slide_idx, "phase": phase, "message": str(d)}


def _ae_violation_to_diag(v) -> dict:
    """Convert AestheticViolation to a flat dict.

    Severity mapping: Verdict.BLOCK → "error", WARN → "warning", ALLOW → "info".
    This is canonical — BLOCK means the build must fail or the agent must explain.
    """
    severity_map = {Verdict.BLOCK: "error", Verdict.WARN: "warning", Verdict.ALLOW: "info"}
    sv = severity_map.get(v.verdict, v.verdict.name.lower() if hasattr(v.verdict, 'name') else str(v.verdict))
    return {
        "kind": v.rule_id,
        "category": v.category,
        "severity": sv,
        "priority": v.priority,
        "elem_id": v.element_id,
        "message": v.message,
        "metrics": v.metrics,
    }


def _render_slide(prs, canvas, template, slide_index=0, total_slides=1):
    """Grid-to-PPT full slide render: background + info layer + decorations + page number."""
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    layout_idx = 0
    try:
        layout_count = len(prs.slide_layouts)
        if layout_count > 0:
            layout_idx = min(slide_index, layout_count - 1)
    except Exception as e:
        print(f"[PPTBuilder] slide_layouts error (template PPTX): {e}")
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Delete all placeholder shapes from the base layout — engine adds its own shapes.
    # Default slide layouts carry title/content placeholders that the grid engine never
    # populates, leaving ghost empty boxes in the output.
    shapes_to_delete = [s for s in slide.shapes if s.is_placeholder]
    for s in shapes_to_delete:
        sp = s._element
        sp.getparent().remove(sp)

    # Background fill
    bg_hex = template.bg_hex.lstrip("#")
    bg_rgb = RGBColor(int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16))
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_rgb

    page_w_pt = canvas.config.canvas_w_pt
    page_h_pt = canvas.config.canvas_h_pt

    # Phase 1: information layer elements
    for pe in canvas._phase1_rects if hasattr(canvas, '_phase1_rects') else {}:
        x, y, w, h = canvas._phase1_rects.get(pe, (0, 0, 0, 0))
        ct, payload = canvas._phase1_payloads.get(pe, (ContentType.UNKNOWN, None))
        if w <= 0 or h <= 0:
            continue

        if ct == ContentType.IMAGE and payload and payload.image_path:
            # Contain-fit image rendering: PIL natural size -> contain -> no crop, no stretch
            _render_image(slide, x, y, w, h, payload)
            # Optional caption
            if payload.caption:
                cap_x = Emu(int(x * 12700))
                cap_y = Emu(int((y + h - 6) * 12700))
                cap_w = Emu(int(w * 12700))
                cap_h = Emu(int(18 * 12700))
                tb = slide.shapes.add_textbox(cap_x, cap_y, cap_w, cap_h)
                tb.text_frame.word_wrap = True
                p = tb.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = payload.caption
                run.font.size = Pt(template.caption_size)
                run.font.color.rgb = RGBColor(*_hex_to_rgb(template.gray_hex) if template.gray_hex else (0x66, 0x66, 0x66))
        else:
            _render_payload(slide, x, y, w, h, ct, payload,
                {"LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER, "RIGHT": PP_ALIGN.RIGHT})

    # Phase 2: decorations
    for dec in (canvas._decoration_payloads if hasattr(canvas, '_decoration_payloads') else []):
        if dec.get("type") == "arrow":
            from pptx.enum.shapes import MSO_CONNECTOR_TYPE
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT,
                Pt(dec["x1"]), Pt(dec["y1"]),
                Pt(dec["x2"]), Pt(dec["y2"]))
            connector.line.color.rgb = RGBColor(*dec.get("line_color", (0x66, 0x66, 0x66)))
            connector.line.width = Pt(dec.get("line_width_pt", 1.5))
            if dec.get("text"):
                tx = (dec["x1"] + dec["x2"]) / 2
                ty = (dec["y1"] + dec["y2"]) / 2
                from pptx.enum.text import MSO_AUTO_SIZE
                label = slide.shapes.add_textbox(Pt(tx - 55), Pt(ty - 16), Pt(110), Pt(32))
                label.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                label.text_frame.word_wrap = True
                label.text_frame.paragraphs[0].text = dec["text"]
                label.text_frame.paragraphs[0].font.size = Pt(dec.get("font_size", 10))
                label.text_frame.paragraphs[0].font.color.rgb = RGBColor(*dec.get("font_color", (0x55, 0x55, 0x55)))

    # Page number
    fn = f"{slide_index + 1}/{total_slides}"
    pn = slide.shapes.add_textbox(Pt(page_w_pt - 60), Pt(page_h_pt - 28), Pt(48), Pt(20))
    pn.text_frame.paragraphs[0].text = fn
    pn.text_frame.paragraphs[0].font.size = Pt(template.page_number_size)
    pn.text_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(template.dim_hex)) if template.dim_hex else RGBColor(0x88, 0x88, 0x99)
    pn.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
