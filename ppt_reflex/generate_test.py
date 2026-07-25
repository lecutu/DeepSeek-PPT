"""
生成测试文件：手工制作包含 15 种问题的 PPT

每页故意加入一种或多种问题：
  1. 元素超出画布
  2. 标题与正文重叠
  3. 两个正文框重叠
  4. 图注遮挡图片关键信息
  5. 页边距不足
  6. 左边缘未对齐
  7. 多元素间距不均
  8. 图片比例失真
  9. 正文字号低于阈值
  10. 文本框疑似溢出
  11. 页脚覆盖内容
  12. 页面内容密度过高
  13. 合法的文字覆盖背景图片 (应不报告)
  14. z-order 导致正文被遮挡
  15. 修复一个问题后引入另一个问题
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path


def make_test_pptx(output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 960 pt = 16:9
    prs.slide_height = Inches(7.5)     # 540 pt
    W = prs.slide_width
    H = prs.slide_height

    SL = prs.slide_layouts[6]  # blank

    def add_slide():
        return prs.slides.add_slide(SL)

    # ═══════════════════════════════════════════════════════
    # Slide 1: 元素超出画布
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    # Title box partially off right edge
    tb = s.shapes.add_textbox(W - Inches(1), Inches(0.5), Inches(3), Inches(1))
    tb.text_frame.text = "超出右边界标题"
    # Body extending beyond bottom
    tb2 = s.shapes.add_textbox(Inches(1), H - Inches(0.3), Inches(4), Inches(2))
    tb2.text_frame.text = "正文超出底部"

    # ═══════════════════════════════════════════════════════
    # Slide 2: 标题与正文重叠
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1.2))
    tb.text_frame.text = "标题"
    tf = tb.text_frame.paragraphs[0]
    tf.runs[0].font.size = Pt(28)
    # Body starts at same y as title ends → deliberate overlap
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(3))
    tb2.text_frame.text = "这段正文与标题区域轻微重叠"
    tf2 = tb2.text_frame.paragraphs[0]
    tf2.runs[0].font.size = Pt(11)  # below threshold

    # ═══════════════════════════════════════════════════════
    # Slide 3: 两个正文框重叠
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(2.5))
    tb.text_frame.text = "第一个正文框（内容 A）"
    tb2 = s.shapes.add_textbox(Inches(3.5), Inches(1.8), Inches(5), Inches(2.5))
    tb2.text_frame.text = "第二个正文框（内容 B）—— 两者大量重叠"

    # ═══════════════════════════════════════════════════════
    # Slide 4: 图注遮挡图片关键信息
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    # Picture placeholder (we use a rectangle as figure stand-in)
    fig = s.shapes.add_shape(1, Inches(1), Inches(0.8), Inches(6), Inches(4))  # rectangle
    fig.fill.solid()
    fig.fill.fore_color.rgb = RGBColor(200, 200, 220)
    fig.text_frame.text = "[图片区域]"
    # Caption overlapping the figure bottom
    cap = s.shapes.add_textbox(Inches(1), Inches(3.5), Inches(6), Inches(2))
    cap.text_frame.text = "图注文字遮挡了图片下半部分"
    # This overlap should be ALLOWED (caption on figure)

    # ═══════════════════════════════════════════════════════
    # Slide 5: 页边距不足 (元素贴近边缘)
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    # Body too close to left edge
    tb = s.shapes.add_textbox(Inches(0.05), Inches(0.5), Inches(5), Inches(3))
    tb.text_frame.text = "正文离左边缘仅 0.05 英寸（约 3.6pt），不足安全边距 36pt"

    # ═══════════════════════════════════════════════════════
    # Slide 6: 左边缘未对齐
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    tb = s.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(5), Inches(0.8))
    tb.text_frame.text = "对齐参考元素"
    tb2 = s.shapes.add_textbox(Inches(1.15), Inches(1.5), Inches(5), Inches(0.8))
    tb2.text_frame.text = "左边缘偏移 0.15 英寸（约 10.8pt）未对齐"
    tb3 = s.shapes.add_textbox(Inches(1.03), Inches(2.5), Inches(5), Inches(0.8))
    tb3.text_frame.text = "轻微偏移 0.03 英寸（约 2.2pt）—— 应自动吸附"

    # ═══════════════════════════════════════════════════════
    # Slide 7: 多元素间距不均
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    # Same widths (1.8" each). Gaps: 2.7, 3.1, 2.52 → mean 2.77
    # dev: 2.5%, 11.8%, 9.1% → too small (< 50%)
    # Let me explicitly make it uneven: position 2 is off
    x_positions = [Inches(0.5), Inches(3.2), Inches(8.5), Inches(10.5)]
    width = Inches(1.8)
    # Gaps from borders: 0.5, 2.7, 3.5, 2.0, 3.3 (out of canvas)
    # Gaps: 2.7, 3.5, 2.0 → mean 2.73, dev: 1%, 28%, 27% → max 28% < 50%
    # Wrong calc. Gaps = x[i+1] - right[i].
    for x in x_positions:
        tb = s.shapes.add_textbox(x, Inches(2), width, Inches(1.5))
        tb.text_frame.text = f"列{len(s.shapes)}"
    # Gaps: 1.95, 3.28, 2.58 in — uneven (mean 2.6, max dev ~26%)

    # ═══════════════════════════════════════════════════════
    # Slide 8: 图片比例失真
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    fig = s.shapes.add_shape(1, Inches(1), Inches(1), Inches(4), Inches(5))
    fig.fill.solid()
    fig.fill.fore_color.rgb = RGBColor(180, 220, 180)
    fig.text_frame.text = "图片比例: 4×5 原始 4:3 → 失真"
    # Original would be 4:3, current is 4:5 = 0.8, deviates significantly

    # ═══════════════════════════════════════════════════════
    # Slide 9: 正文字号低于阈值
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(0.8))
    tb.text_frame.text = "标题（正常 28pt）"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(4))
    tb2.text_frame.text = "正文 10pt 太小，低于 14pt 下限"
    tb2.text_frame.paragraphs[0].runs[0].font.size = Pt(10)

    # ═══════════════════════════════════════════════════════
    # Slide 10: 文本框疑似溢出（窄框+长文字）
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1.5), Inches(0.6))
    tb.text_frame.text = "这段文字非常非常多放不下"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    # Small box, large text at 14pt — will likely overflow

    # ═══════════════════════════════════════════════════════
    # Slide 11: 页脚覆盖内容
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(5.0))
    tb.text_frame.text = "正文内容延伸到页面底部"
    # Footer intentionally overlapping the body text
    footer = s.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(6), Inches(2.0))
    footer.text_frame.text = "页脚文字覆盖了正文底部区域"

    # ═══════════════════════════════════════════════════════
    # Slide 12: 页面内容密度过高
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    # Jam 12 text boxes, kept within safe margins
    for row in range(3):
        for col in range(4):
            x = Inches(0.5 + col * 3.1)
            y = Inches(0.5 + row * 2.2)
            tb = s.shapes.add_textbox(x, y, Inches(2.8), Inches(1.8))
            tb.text_frame.text = f"R{row}C{col} 内容"
            tb.text_frame.paragraphs[0].runs[0].font.size = Pt(9)

    # ═══════════════════════════════════════════════════════
    # Slide 13: 合法重叠——文字覆盖背景图片（不应报告）
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    bg = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(50, 50, 100)
    bg.text_frame.text = "背景"
    # Text intentionally over background — this is legitimate
    tb = s.shapes.add_textbox(Inches(2), Inches(2), Inches(8), Inches(3))
    tb.text_frame.text = "覆盖在背景上的文字——这是有意的设计"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    # We'll mark the bg as "background" role manually

    # ═══════════════════════════════════════════════════════
    # Slide 14: z-order 导致正文被遮挡
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    body = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
    body.text_frame.text = "这段正文应该可见但被遮挡"
    # Blocking shape added AFTER, sits on top
    block = s.shapes.add_shape(1, Inches(2), Inches(1.5), Inches(4), Inches(2))
    block.fill.solid()
    block.fill.fore_color.rgb = RGBColor(255, 200, 200)
    block.text_frame.text = "遮挡块"

    # ═══════════════════════════════════════════════════════
    # Slide 15: 修复一个问题后引入另一个问题
    # ═══════════════════════════════════════════════════════
    s = add_slide()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(2))
    tb.text_frame.text = "正文 A"
    tb2 = s.shapes.add_textbox(Inches(2), Inches(1.5), Inches(5), Inches(2))
    tb2.text_frame.text = "正文 B（与 A 重叠）"
    tb3 = s.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(11), Inches(1))
    tb3.text_frame.text = "如果把 B 右移，它可能越界"
    # A and B overlap → need separation. But B is near right edge.
    # Moving B right to resolve collision may push it out of bounds.

    # ── Save ───────────────────────────────────────────────
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Created test file: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    target = os.path.join(os.path.dirname(os.path.abspath(__file__)), "cases", "broken.pptx")
    make_test_pptx(target)
