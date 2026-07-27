"""
ppt_reflex/roundtrip_check.py — Reopen saved PPTX and verify every text box against engine estimates.

This is the ground-truth feedback loop: estimate → render → read back → detect gaps.
python-pptx font metrics differ from engine heuristics; only the saved file tells the truth.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Pt
from .grid.text_metrics import estimate_text_size, _line_width


def check_overflow(path: str) -> list[dict]:
    """Reopen a saved PPTX and verify every text shape for overflow — 2D (vertical + horizontal).

    Returns a list of overflow diagnostics per shape, sorted by severity:
      - silent_overflow: shape height < estimated height by >4pt → error
      - overflow_horizontal: longest non-breakable word > box width → error
      - tight_fit: shape height within 2–4pt of estimate → warning
      - ok: everything else

    Shape types checked: TEXT_BOX, AUTO_SHAPE with text, PLACEHOLDER with text.
    """
    prs = Presentation(path)
    results: list[dict] = []

    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text.strip()
            if not text:
                continue

            actual_w = shape.width / 12700
            actual_h = shape.height / 12700

            # Get dominant font size across all runs (mode, not just first)
            font_sizes: dict[float, int] = {}
            line_spacing = 1.2
            for para in tf.paragraphs:
                if para.line_spacing:
                    line_spacing = para.line_spacing / 12700 / max(12.0, 0.1)
                for run in para.runs:
                    if run.font.size:
                        fs = run.font.size / 12700
                        font_sizes[fs] = font_sizes.get(fs, 0) + 1
            font_size = max(font_sizes, key=font_sizes.get) if font_sizes else 12.0

            # ── Vertical overflow ──
            _, ov_y, rw, rh = estimate_text_size(
                text, font_pt=font_size, line_spacing=line_spacing,
                box_width_pt=actual_w, box_height_pt=9999, word_wrap=True,
            )

            gap_v = rh - actual_h
            if gap_v > 4:
                severity = "error"
                kind = "silent_overflow"
            elif gap_v > 2:
                severity = "warning"
                kind = "tight_fit"
            else:
                severity = "ok"
                kind = "ok"

            if severity != "ok":
                results.append({
                    "slide": si,
                    "kind": kind,
                    "severity": severity,
                    "shape_type": str(shape.shape_type),
                    "position": (
                        round(shape.left / 12700, 0),
                        round(shape.top / 12700, 0),
                    ),
                    "actual_size": (round(actual_w, 0), round(actual_h, 0)),
                    "estimated_h": round(rh, 0),
                    "gap_pt": round(gap_v, 1),
                    "font_size": round(font_size, 1),
                    "auto_size": str(tf.auto_size),
                    "text_preview": text[:80],
                    "message": (
                        f"text needs {rh:.0f}pt, box is {actual_h:.0f}pt — "
                        f"{gap_v:.0f}pt hidden below. font={font_size:.0f}pt, auto_size={tf.auto_size}"
                    ),
                })

            # ── Horizontal overflow: longest unbreakable word vs box width ──
            longest_word_w = 0.0
            for line in text.split("\n"):
                for word in line.split(" "):
                    ww = _line_width(word, font_size)
                    if ww > longest_word_w:
                        longest_word_w = ww
            if actual_w > 0 and longest_word_w > actual_w + 1:
                results.append({
                    "slide": si,
                    "kind": "overflow_horizontal",
                    "severity": "error",
                    "shape_type": str(shape.shape_type),
                    "position": (
                        round(shape.left / 12700, 0),
                        round(shape.top / 12700, 0),
                    ),
                    "actual_size": (round(actual_w, 0), round(actual_h, 0)),
                    "longest_word_w": round(longest_word_w, 1),
                    "box_w": round(actual_w, 0),
                    "font_size": round(font_size, 1),
                    "auto_size": str(tf.auto_size),
                    "text_preview": text[:80],
                    "message": (
                        f"longest word {longest_word_w:.1f}pt > box {actual_w:.0f}pt — "
                        f"horizontal overflow. font={font_size:.0f}pt"
                    ),
                })

    results.sort(key=lambda d: {"error": 0, "warning": 1, "ok": 2}[d["severity"]])
    return results
